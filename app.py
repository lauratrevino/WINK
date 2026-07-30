import os, json, uuid, secrets, traceback, time, threading
from collections import defaultdict, deque
from datetime import datetime, timedelta
from flask import (Flask, render_template, request, jsonify,
                   session, redirect, url_for, g)
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash

app = Flask(__name__)
_secret = os.environ.get("SECRET_KEY")
if not _secret:
    # A hardcoded fallback secret key lets anyone forge session cookies
    # (including an "is admin" session) since Flask signs sessions with this
    # value. Generate a random one instead so a missing env var fails safe —
    # sessions just won't survive a restart until SECRET_KEY is actually set.
    _secret = secrets.token_hex(32)
    print("WARNING: SECRET_KEY not set — using a random per-process key. "
          "Sessions will be invalidated on every restart until you set "
          "SECRET_KEY in the environment.")
app.secret_key = _secret
app.config.update(
    SESSION_COOKIE_HTTPONLY=True,
    SESSION_COOKIE_SECURE=os.environ.get("FLASK_ENV") != "development",
    SESSION_COOKIE_SAMESITE="Lax",
    PERMANENT_SESSION_LIFETIME=timedelta(days=7),
)

BASE_DIR      = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")
ALLOWED_EXT   = {"pdf","docx","txt","pptx","xlsx","png","jpg","jpeg"}
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config["MAX_CONTENT_LENGTH"] = 16 * 1024 * 1024

DB_URL            = os.environ.get("DATABASE_URL", "")
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
ADMIN_EMAIL = os.environ.get("ADMIN_EMAIL", "lhall@utep.edu").lower()
MAX_DOCS_PER_STUDENT = 20
# Password reset has no real email provider wired up yet (see forgot_password
# below). Never expose the raw reset link in an HTTP response in production —
# that turns "forgot password" into a way to take over ANY account just by
# knowing their email. Only set this to true for local/dev testing.
DEBUG_SHOW_RESET_LINKS = os.environ.get("DEBUG_SHOW_RESET_LINKS", "false").lower() == "true"

# ── Email ─────────────────────────────────────────────────────
# Standard SMTP config. Works with SendGrid, Mailgun, Amazon SES (SMTP
# interface), Gmail (with an app password), etc. — set these four env vars
# and password reset + deadline reminder emails start actually sending
# instead of only being logged.
SMTP_HOST = os.environ.get("SMTP_HOST", "")
SMTP_PORT = int(os.environ.get("SMTP_PORT", "587"))
SMTP_USER = os.environ.get("SMTP_USER", "")
SMTP_PASS = os.environ.get("SMTP_PASS", "")
FROM_EMAIL = os.environ.get("FROM_EMAIL", SMTP_USER or "wink@utep.edu")
EMAIL_CONFIGURED = bool(SMTP_HOST and SMTP_USER and SMTP_PASS)
# Shared secret that an external scheduler (Render cron job, GitHub Action,
# etc.) must pass to trigger deadline reminder emails, so the endpoint can't
# be used by a random visitor to spam every student.
CRON_SECRET = os.environ.get("CRON_SECRET", "")

def send_email(to_email, subject, body):
    """Send a plain-text email. Returns True on success. Falls back to
    logging (never raises) if SMTP isn't configured or sending fails, so a
    flaky email provider never breaks the request that triggered it."""
    if not EMAIL_CONFIGURED:
        print(f"EMAIL (not sent — SMTP not configured) to={to_email} subject={subject!r}")
        return False
    try:
        import smtplib
        from email.mime.text import MIMEText
        msg = MIMEText(body)
        msg["Subject"] = subject
        msg["From"] = FROM_EMAIL
        msg["To"] = to_email
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT, timeout=10) as server:
            server.starttls()
            server.login(SMTP_USER, SMTP_PASS)
            server.sendmail(FROM_EMAIL, [to_email], msg.as_string())
        return True
    except Exception as e:
        print(f"send_email error to={to_email}: {e}")
        return False

# ── Cost controls ────────────────────────────────────────────
# Model: Haiku is ~3x cheaper than Sonnet on both input and output tokens.
# Override with the WINK_MODEL env var if you want to trade cost for quality.
CHAT_MODEL = os.environ.get("WINK_MODEL", "claude-haiku-4-5-20251001")
CHAT_MAX_TOKENS = 1024
# Hard cap on how many characters of document text get sent to the model per
# question, regardless of how many documents the student has uploaded. This
# is the single biggest cost lever: every uploaded page otherwise gets resent
# on every single question, with no caching, for the life of the conversation.
# Raised from 24K: build_doc_context() now divides this budget evenly across
# every uploaded document instead of first-come-first-served, so a slightly
# bigger pool means each individual document still gets a workable amount of
# its own content rather than being squeezed to almost nothing.
MAX_DOC_CONTEXT_CHARS = 40000
# Separate budget for admin-uploaded "general" reference documents that apply
# to every student (see build_global_doc_context()) — kept apart from the
# per-student budget above so the two caches don't invalidate each other.
MAX_GLOBAL_DOC_CONTEXT_CHARS = 20000
# Deadline extraction runs ONCE per upload, not on every question — it should
# NOT share the tight per-message chat budget above. Match extract_text()'s
# own 60,000-char storage cap instead, so extraction sees everything that was
# actually kept from the document (was previously truncated to the same
# ~24K used for live chat, which could cut off a syllabus's schedule section
# entirely on longer documents — the actual bug behind incomplete deadlines).
DEADLINE_EXTRACTION_MAX_CHARS = 60000
# Cap for a "temporary, this-conversation-only" upload (see /upload's
# `temporary` flag). These are never written to the documents table or
# counted against MAX_DOCS_PER_STUDENT — the extracted text is handed back
# to the client, which resends it with each /chat call in that conversation
# only. Kept modest since it's on top of the student's regular doc context.
MAX_TEMP_DOC_CHARS = 20000
# How many prior chat messages (user+assistant turns) to actually send with
# each request. Conversation history otherwise grows unbounded and gets
# re-billed as input tokens on every new question.
MAX_CHAT_HISTORY_MESSAGES = 12
# Cap how many web searches Claude can run per question (each search is
# $0.01 regardless of whether Claude uses the results).
WEB_SEARCH_MAX_USES = 3
# Reject absurdly long single messages instead of billing (and paying) for
# whatever a script or a copy-pasted textbook chapter throws at the endpoint.
MAX_USER_MESSAGE_CHARS = 6000

# ── Abuse protection ─────────────────────────────────────────
# Best-effort, per-process rate limiting. This resets if the process
# restarts and isn't shared across gunicorn workers, so it's not a hard
# guarantee — but it's a cheap, dependency-free way to blunt brute-force
# login attempts and runaway/scripted chat spam. For a hard guarantee under
# multiple workers/instances, back this with Redis instead.
_rate_lock = threading.Lock()
_rate_buckets = defaultdict(deque)

def rate_limited(key, max_calls, window_seconds):
    """Returns 0 if the call is allowed, otherwise the number of seconds
    until the oldest call in the window ages out (so the caller can tell a
    client exactly how long to back off, rather than just "try later").
    Every existing call site does `if rate_limited(...):`, which still works
    unchanged — 0 is falsy, any positive number of seconds is truthy."""
    now = time.time()
    with _rate_lock:
        bucket = _rate_buckets[key]
        while bucket and now - bucket[0] > window_seconds:
            bucket.popleft()
        if len(bucket) >= max_calls:
            return round(window_seconds - (now - bucket[0]), 1)
        bucket.append(now)
        return 0

# ── Speed controls ───────────────────────────────────────────
# Build the Anthropic client once, at import time, and reuse it for every
# request. Creating a fresh httpx.Client per request (the old behavior) means
# a brand-new TCP + TLS handshake to Anthropic's servers on every single
# question. Reusing one client with a connection pool lets requests reuse an
# already-open, already-authenticated connection — this alone typically saves
# 100-300ms of pure connection setup time before the model even starts
# thinking.
import httpx, anthropic as ac
_http_client = httpx.Client(
    timeout=httpx.Timeout(110.0, connect=5.0),
    limits=httpx.Limits(max_keepalive_connections=20, max_connections=50),
    http2=True,
)
anthropic_client = ac.Anthropic(api_key=ANTHROPIC_API_KEY, http_client=_http_client) if ANTHROPIC_API_KEY else None

CLASSIFICATIONS = ["Freshman","Sophomore","Junior","Senior","Graduate","Faculty"]
MAJORS = [
    "Accounting","Biology","Business Administration","Chemistry",
    "Civil Engineering","Communication","Computer Science",
    "Criminal Justice","Economics","Education","Electrical Engineering",
    "English","Environmental Science","Finance","History",
    "Industrial Engineering","Information Systems","Kinesiology",
    "Management","Marketing","Mathematics","Mechanical Engineering",
    "Nursing","Political Science","Psychology","Public Health",
    "Social Work","Sociology","Spanish","Other"
]

# ── Text Extraction ───────────────────────────────────────
def extract_text(filepath, orig_name):
    ext = orig_name.rsplit(".", 1)[-1].lower() if "." in orig_name else ""
    text = ""
    try:
        if ext == "txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == "pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(filepath)
                pages = []
                for i, page in enumerate(reader.pages):
                    try:
                        t = page.extract_text()
                        if t and t.strip():
                            pages.append(f"[Page {i+1}]\n{t.strip()}")
                    except Exception as pe:
                        print(f"PDF page {i+1} error: {pe}")
                text = "\n\n".join(pages)
                print(f"PDF extracted {len(text)} chars from {len(reader.pages)} pages")
            except Exception as e:
                print(f"PDF extract failed: {e}"); traceback.print_exc()
        elif ext == "docx":
            try:
                from docx import Document
                doc = Document(filepath)
                parts = []
                for p in doc.paragraphs:
                    if p.text.strip():
                        parts.append(p.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                text = "\n".join(parts)
                print(f"DOCX extracted {len(text)} chars")
            except Exception as e:
                print(f"DOCX extract failed: {e}"); traceback.print_exc()
        elif ext == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                slides = []
                for i, slide in enumerate(prs.slides):
                    parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text.strip())
                    if parts:
                        slides.append(f"[Slide {i+1}]\n" + "\n".join(parts))
                text = "\n\n".join(slides)
                print(f"PPTX extracted {len(text)} chars")
            except Exception as e:
                print(f"PPTX extract failed: {e}"); traceback.print_exc()
        elif ext == "xlsx":
            try:
                from openpyxl import load_workbook
                # read_only + data_only: stream rows instead of loading the
                # whole workbook into memory, and read formula *results*
                # rather than the formula text itself
                wb = load_workbook(filepath, read_only=True, data_only=True)
                sheets = []
                for ws in wb.worksheets:
                    lines = []
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                        if cells:
                            lines.append(" | ".join(cells))
                        if len(lines) >= 2000:  # guard against extreme sheets
                            lines.append("[...sheet truncated...]")
                            break
                    if lines:
                        sheets.append(f"[Sheet: {ws.title}]\n" + "\n".join(lines))
                n_sheets = len(wb.worksheets)
                wb.close()
                text = "\n\n".join(sheets)
                print(f"XLSX extracted {len(text)} chars from {n_sheets} sheet(s)")
            except Exception as e:
                print(f"XLSX extract failed: {e}"); traceback.print_exc()
        elif ext in ("jpg","jpeg","png"):
            text = f"[Image file: {orig_name}]"
        else:
            try:
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            except:
                text = ""
    except Exception as e:
        print(f"extract_text error for {orig_name}: {e}"); text = ""
    if len(text) > 60000:
        text = text[:60000] + "\n\n[Document truncated at 60,000 characters]"
    return text.strip()

def extract_deadlines(content, today=None):
    """Ask Claude to pull structured (title, due_date) pairs out of a
    document's text — e.g. a syllabus's assignment schedule. Returns a list
    of {"title": str, "due_date": "YYYY-MM-DD"} dicts, or [] on any failure
    (no document content, model error, unparsable response, etc). This is a
    small, cheap, one-time Haiku call made once per upload, not per question.
    """
    if not anthropic_client or not content or not content.strip():
        return []
    today = today or datetime.utcnow().strftime("%Y-%m-%d")
    try:
        resp = anthropic_client.messages.create(
            model=CHAT_MODEL,
            max_tokens=800,
            system=(
                "Extract assignment, exam, and other academic deadlines from the "
                "document text the user provides. Respond with ONLY a JSON array "
                "(no prose, no markdown fences) of objects shaped like "
                '{"title": "...", "due_date": "YYYY-MM-DD"}. '
                f"Today's date is {today} — resolve relative or partial dates "
                "(e.g. \"March 3\" with no year, or \"next Friday\") against it. "
                "Skip anything without a specific date you can resolve. "
                "If there are no clear deadlines, respond with []."
            ),
            messages=[{"role": "user", "content": content[:DEADLINE_EXTRACTION_MAX_CHARS]}],
        )
        raw = "".join(b.text for b in resp.content if getattr(b, "type", None) == "text").strip()
        if raw.startswith("```"):
            raw = raw.strip("`").split("\n", 1)[-1] if "\n" in raw else raw.strip("`")
        items = json.loads(raw)
        out = []
        for it in items if isinstance(items, list) else []:
            title = str(it.get("title","")).strip()[:200]
            due   = str(it.get("due_date","")).strip()
            try:
                datetime.strptime(due, "%Y-%m-%d")
            except Exception:
                continue
            if title:
                out.append({"title": title, "due_date": due})
        return out[:30]
    except Exception as e:
        print(f"extract_deadlines error: {e}")
        return []

def build_doc_context(docs):
    if not docs:
        return "\n\nThe student has not uploaded any course documents yet."
    has_content = any((d.get("content") or "").strip() for d in docs)
    if not has_content:
        ctx = f"\n\nThe student has {len(docs)} uploaded file(s) but no text could be extracted. "
        ctx += "Files: " + ", ".join(d["orig_name"] for d in docs)
        return ctx
    ctx = f"\n\n{'='*60}\nSTUDENT'S UPLOADED COURSE DOCUMENTS ({len(docs)} files)\n"
    ctx += ("Every document the student has uploaded is listed below in full or in part — none "
            "have been skipped. Never tell the student to re-upload something that appears here; "
            "if you only see part of a long one, say you have it but only part of its content, "
            "and offer to look at a specific section if they ask.\n")
    ctx += "Answer questions using the actual content of these documents.\n"
    ctx += "Quote specific text, deadlines, requirements directly from the documents.\n"
    ctx += f"{'='*60}\n\n"
    # Give every document a fair, even share of the total budget up front — a
    # first-come-first-served allocation could let one or two large, recently
    # uploaded files consume the whole budget and push older documents out
    # entirely, which is exactly the "it won't read documents I already
    # uploaded" bug. Every document is now guaranteed at least some content.
    per_doc_budget = max(MAX_DOC_CONTEXT_CHARS // max(len(docs), 1), 2000)
    for i, d in enumerate(docs):
        content = (d.get("content") or "").strip()
        header = f"[DOCUMENT {i+1}] {d['orig_name']}\n"
        header += f"Course: {d['course']} | Size: {round(d.get('size_bytes',0)/1024,1)} KB\n"
        header += f"Content ({len(content)} chars):\n"
        if len(content) > per_doc_budget:
            content = content[:per_doc_budget] + "\n[Shortened here to fit — ask about this document specifically for more of it.]"
        ctx += header
        ctx += content if content else "[No text could be extracted from this file]"
        ctx += f"\n\n{'-'*40}\n\n"
    ctx += f"{'='*60}\n"
    return ctx

def get_global_docs(university=None):
    """Institution-wide reference documents visible to every student's chat
    context (for their own university only) but never shown in any student's
    own document list or counted against their upload cap — stored as
    documents with student_id NULL, tagged with the target university."""
    if not DB_URL: return []
    try:
        conn = get_db(); cur = conn.cursor()
        if university:
            cur.execute("""SELECT * FROM documents WHERE student_id IS NULL
                           AND lower(university)=lower(%s) ORDER BY uploaded_at DESC""",
                        (university,))
        else:
            cur.execute("SELECT * FROM documents WHERE student_id IS NULL ORDER BY uploaded_at DESC")
        docs = [dict(r) for r in cur.fetchall()]; cur.close(); db_release(conn)
        return docs
    except Exception as e:
        print(f"get_global_docs error: {e}"); return []

def build_global_doc_context(docs, university=None):
    if not docs:
        return "\n\nNo general reference documents have been added yet."
    label = f" for {university}" if university else ""
    ctx = f"\n\n{'='*60}\nGENERAL REFERENCE DOCUMENTS (apply to every student{label}, not just this one)\n"
    ctx += ("These were uploaded by an administrator and are not visible to the student as their "
            "own files — don't refer to them as 'your uploaded documents' or mention that they "
            "were uploaded separately; just use them as background knowledge when relevant.\n")
    ctx += f"{'='*60}\n\n"
    per_doc_budget = max(MAX_GLOBAL_DOC_CONTEXT_CHARS // max(len(docs), 1), 2000)
    for i, d in enumerate(docs):
        content = (d.get("content") or "").strip()
        if not content:
            continue
        if len(content) > per_doc_budget:
            content = content[:per_doc_budget] + "\n[Shortened here to fit.]"
        ctx += f"[REFERENCE {i+1}] {d['orig_name']} ({d.get('course') or 'General'})\n{content}\n\n{'-'*40}\n\n"
    ctx += f"{'='*60}\n"
    return ctx

# ── DB ────────────────────────────────────────────────────
# Speed: pool connections instead of opening a brand-new TCP + auth
# handshake to Postgres on every single query. A typical request makes
# several DB round trips (auth check, doc lookup, event log, ...); pooling
# turns most of those into "grab an already-open connection" instead of
# "negotiate a new one", which matters a lot under concurrent load.
import psycopg2
from psycopg2 import pool as _pg_pool
from psycopg2.extras import RealDictCursor

_db_pool = None
if DB_URL:
    try:
        _db_pool = _pg_pool.ThreadedConnectionPool(1, 20, DB_URL, cursor_factory=RealDictCursor)
    except Exception as e:
        print(f"DB pool init failed, falling back to per-request connections: {e}")
        _db_pool = None

def get_db():
    """Returns one connection per request, cached on Flask's `g`. The first
    call in a request checks a connection out of the pool; every subsequent
    call in that same request (current_student(), get_docs(), log_event(),
    the route's own queries, ...) reuses it instead of checking out another.
    This also means release no longer depends on every call site remembering
    to call db_release() on every code path — see teardown handler below,
    which guarantees it runs exactly once per request even after an
    exception. (db_release() itself is now a no-op — kept so the ~30 existing
    call sites don't need to change.)
    """
    if getattr(g, "_db_conn", None) is not None:
        return g._db_conn
    if _db_pool:
        conn = _db_pool.getconn()
    else:
        conn = psycopg2.connect(DB_URL, cursor_factory=RealDictCursor)
    g._db_conn = conn
    return conn

def db_release(conn):
    """No-op: kept for source compatibility with existing call sites.
    Actual release happens once per request in the teardown handler below —
    releasing here too would return the same pooled connection twice (once
    now, once at teardown), which could hand it to two requests at once."""
    pass

@app.teardown_appcontext
def _release_request_db_connection(exception=None):
    """Guarantees the per-request connection (see get_db() above) always
    goes back to the pool exactly once, whether the request succeeded,
    returned an error response, or raised an uncaught exception. This is
    what actually fixes connection leaks — previously, an exception between
    `get_db()` and the function's own `db_release(conn)` call meant that
    connection was gone from the pool for good."""
    conn = g.pop("_db_conn", None)
    if conn is None:
        return
    if _db_pool:
        try:
            _db_pool.putconn(conn)
        except Exception:
            pass
    else:
        try:
            conn.close()
        except Exception:
            pass

def init_db():
    if not DB_URL:
        print("WARNING: No DATABASE_URL set.")
        return
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("""CREATE TABLE IF NOT EXISTS students (
            id SERIAL PRIMARY KEY, email TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL, first_name TEXT NOT NULL,
            last_name TEXT NOT NULL, classification TEXT NOT NULL,
            major TEXT NOT NULL, university TEXT DEFAULT '',
            created_at TIMESTAMP DEFAULT NOW())""")
        cur.execute("ALTER TABLE students ADD COLUMN IF NOT EXISTS university TEXT DEFAULT ''")
        cur.execute("ALTER TABLE students ADD COLUMN IF NOT EXISTS is_active BOOLEAN DEFAULT TRUE")
        cur.execute("ALTER TABLE students ADD COLUMN IF NOT EXISTS email_verified BOOLEAN DEFAULT FALSE")
        cur.execute("ALTER TABLE students ADD COLUMN IF NOT EXISTS verification_token TEXT")
        cur.execute("""CREATE TABLE IF NOT EXISTS documents (
            id SERIAL PRIMARY KEY,
            student_id INTEGER REFERENCES students(id) ON DELETE CASCADE,
            filename TEXT NOT NULL, orig_name TEXT NOT NULL,
            course TEXT NOT NULL, size_bytes INTEGER DEFAULT 0,
            content TEXT DEFAULT '',
            uploaded_at TIMESTAMP DEFAULT NOW())""")
        cur.execute("ALTER TABLE documents ADD COLUMN IF NOT EXISTS content TEXT DEFAULT ''")
        cur.execute("ALTER TABLE documents ADD COLUMN IF NOT EXISTS crn TEXT DEFAULT ''")
        # Only meaningful on global/admin-uploaded rows (student_id IS NULL) —
        # scopes each general-reference document to one school's knowledge
        # base, now that WINK serves more than just UTEP.
        cur.execute("ALTER TABLE documents ADD COLUMN IF NOT EXISTS university TEXT DEFAULT ''")
        cur.execute("CREATE INDEX IF NOT EXISTS idx_documents_global_university ON documents(university) WHERE student_id IS NULL")
        # Use TEXT for payload — simple and reliable across all Postgres versions
        cur.execute("""CREATE TABLE IF NOT EXISTS events (
            id SERIAL PRIMARY KEY, student_id INTEGER,
            event_type TEXT NOT NULL,
            payload TEXT DEFAULT '{}',
            created_at TIMESTAMP DEFAULT NOW())""")
        cur.execute("""CREATE TABLE IF NOT EXISTS password_resets (
            id SERIAL PRIMARY KEY,
            student_id INTEGER REFERENCES students(id) ON DELETE CASCADE,
            token TEXT UNIQUE NOT NULL,
            expires_at TIMESTAMP NOT NULL,
            used BOOLEAN DEFAULT FALSE,
            created_at TIMESTAMP DEFAULT NOW())""")
        # Deadlines extracted from uploaded documents (syllabi, assignment sheets)
        cur.execute("""CREATE TABLE IF NOT EXISTS deadlines (
            id SERIAL PRIMARY KEY,
            student_id INTEGER REFERENCES students(id) ON DELETE CASCADE,
            document_id INTEGER REFERENCES documents(id) ON DELETE CASCADE,
            course TEXT NOT NULL,
            title TEXT NOT NULL,
            due_date DATE,
            reminded BOOLEAN DEFAULT FALSE,
            created_at TIMESTAMP DEFAULT NOW())""")
        # Saved chat conversations, so students can revisit and export past sessions
        cur.execute("""CREATE TABLE IF NOT EXISTS conversations (
            id SERIAL PRIMARY KEY,
            student_id INTEGER REFERENCES students(id) ON DELETE CASCADE,
            title TEXT DEFAULT 'New conversation',
            messages TEXT DEFAULT '[]',
            share_token TEXT UNIQUE,
            created_at TIMESTAMP DEFAULT NOW(),
            updated_at TIMESTAMP DEFAULT NOW())""")
        # Foreign keys in Postgres do NOT automatically index the referencing
        # column, and these are exactly the columns get_docs(), log_event(),
        # and the analytics queries filter on for every single request. Without
        # these, those queries full-table-scan documents/events as they grow —
        # slower and more expensive DB CPU with every new student and question.
        cur.execute("CREATE INDEX IF NOT EXISTS idx_documents_student_id ON documents(student_id)")
        cur.execute("CREATE INDEX IF NOT EXISTS idx_events_student_type ON events(student_id, event_type)")
        cur.execute("CREATE INDEX IF NOT EXISTS idx_events_type ON events(event_type)")
        cur.execute("CREATE INDEX IF NOT EXISTS idx_events_created_at ON events(created_at)")
        cur.execute("CREATE INDEX IF NOT EXISTS idx_deadlines_student_id ON deadlines(student_id)")
        cur.execute("CREATE INDEX IF NOT EXISTS idx_deadlines_due_date ON deadlines(due_date)")
        cur.execute("CREATE INDEX IF NOT EXISTS idx_conversations_student_id ON conversations(student_id)")
        conn.commit(); cur.close(); db_release(conn)
        print("DB initialized OK.")
    except Exception as e:
        print(f"DB init error: {e}"); traceback.print_exc()

with app.app_context():
    init_db()

# ── Helpers ───────────────────────────────────────────────
def current_student():
    if "sid" not in session or not DB_URL:
        return None
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT * FROM students WHERE id=%s", (session["sid"],))
        s = cur.fetchone(); cur.close(); db_release(conn)
        if s and not s.get("is_active", True):
            session.clear()
            return None
        return dict(s) if s else None
    except Exception as e:
        print(f"current_student error: {e}"); return None

def log_event(sid, etype, payload=None):
    """Log every user action to the events table."""
    if not DB_URL:
        return
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute(
            "INSERT INTO events(student_id, event_type, payload) VALUES(%s, %s, %s)",
            (sid, etype, json.dumps(payload or {}))
        )
        conn.commit(); cur.close(); db_release(conn)
    except Exception as e:
        print(f"log_event ERROR: {e}")
        traceback.print_exc()

def get_docs(sid):
    if not DB_URL: return []
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT * FROM documents WHERE student_id=%s ORDER BY uploaded_at DESC", (sid,))
        docs = [dict(r) for r in cur.fetchall()]; cur.close(); db_release(conn)
        return docs
    except Exception as e:
        print(f"get_docs error: {e}"); return []

def group_docs_by_course(docs):
    """Group documents by course (and CRN, when present)."""
    groups = {}
    order = []
    for d in docs:
        course = (d.get("course") or "").strip() or "General"
        crn = (d.get("crn") or "").strip()
        label = f"{course} (CRN {crn})" if crn else course
        if label not in groups:
            groups[label] = []
            order.append(label)
        groups[label].append(d)
    order.sort()
    return [(label, groups[label]) for label in order]

def safe_payload(raw):
    """Safely parse a payload value regardless of whether it's str, dict, or None."""
    if raw is None:
        return {}
    if isinstance(raw, dict):
        return raw
    try:
        return json.loads(raw)
    except:
        return {}

def compute_engagement_insights(cur):
    """Everything beyond the basic counts already in /analytics-data-full:
    per-university breakdown, session duration, retention, time-to-first-question,
    peak usage heatmap, deadline-driven spikes, upload mix, stale docs, and a
    rough 'general reference material was available' rate. All computed from
    the existing students/events/documents/deadlines tables — no schema or
    frontend instrumentation changes required."""
    out = {}

    # ── Per-university breakdown ──
    cur.execute("""
        SELECT COALESCE(NULLIF(s.university,''), 'Not set') as university,
               COUNT(DISTINCT s.id) as students,
               COUNT(*) FILTER (WHERE e.event_type IN ('login','account_created')) as sessions,
               COUNT(*) FILTER (WHERE e.event_type='question_asked') as questions,
               COUNT(*) FILTER (WHERE e.event_type='file_uploaded') as uploads
        FROM students s LEFT JOIN events e ON e.student_id = s.id
        GROUP BY 1 ORDER BY students DESC""")
    out["by_university"] = [dict(r) for r in cur.fetchall()]

    # ── Session duration (derived from event timestamps — a new session
    #    starts at each login/account_created; its duration is the gap to the
    #    last event before the next session starts) ──
    cur.execute("""
        SELECT e.student_id, e.event_type, e.created_at, COALESCE(NULLIF(s.university,''),'Not set') as university
        FROM events e JOIN students s ON s.id = e.student_id
        ORDER BY e.student_id, e.created_at ASC""")
    rows = cur.fetchall()
    sessions_by_student = {}
    cur_session = None
    for r in rows:
        sid = r["student_id"]
        if sid not in sessions_by_student:
            sessions_by_student[sid] = []
        if r["event_type"] in ("login", "account_created"):
            cur_session = {"university": r["university"], "start": r["created_at"], "end": r["created_at"]}
            sessions_by_student[sid].append(cur_session)
        elif sessions_by_student[sid]:
            sessions_by_student[sid][-1]["end"] = r["created_at"]

    all_durations = []
    durations_by_university = {}
    for sid, sess_list in sessions_by_student.items():
        for sess in sess_list:
            mins = (sess["end"] - sess["start"]).total_seconds() / 60.0
            mins = max(0.0, min(mins, 240.0))  # ignore/cap runaway gaps
            all_durations.append(mins)
            durations_by_university.setdefault(sess["university"], []).append(mins)

    out["avg_session_minutes"] = round(sum(all_durations) / len(all_durations), 1) if all_durations else 0
    out["avg_session_minutes_by_university"] = {
        u: round(sum(v) / len(v), 1) for u, v in durations_by_university.items()
    }

    # ── Retention: % of students active across 2+ distinct weeks ──
    cur.execute("""
        SELECT student_id, COUNT(DISTINCT date_trunc('week', created_at)) as weeks
        FROM events GROUP BY student_id""")
    week_rows = cur.fetchall()
    total_active = len(week_rows)
    returning = sum(1 for r in week_rows if r["weeks"] >= 2)
    out["retention_pct"] = round(returning / total_active * 100, 1) if total_active else 0

    # ── Time-to-first-question ──
    cur.execute("""
        SELECT s.created_at as joined, MIN(e.created_at) as first_q
        FROM students s JOIN events e ON e.student_id = s.id AND e.event_type = 'question_asked'
        GROUP BY s.id, s.created_at""")
    gaps = [(r["first_q"] - r["joined"]).total_seconds() / 60.0 for r in cur.fetchall()]
    gaps = [g for g in gaps if g >= 0]
    out["avg_minutes_to_first_question"] = round(sum(gaps) / len(gaps), 1) if gaps else None

    # ── Peak usage heatmap (questions asked, by day-of-week x hour) ──
    cur.execute("""
        SELECT EXTRACT(DOW FROM created_at)::int as dow, EXTRACT(HOUR FROM created_at)::int as hour, COUNT(*) as n
        FROM events WHERE event_type='question_asked' GROUP BY 1,2""")
    grid = [[0]*24 for _ in range(7)]
    for r in cur.fetchall():
        grid[r["dow"]][r["hour"]] = r["n"]
    out["usage_heatmap"] = grid

    # ── Deadline-driven spikes: questions asked on/around days with deadlines due ──
    cur.execute("SELECT due_date, COUNT(*) as n FROM deadlines WHERE due_date IS NOT NULL GROUP BY due_date")
    due_by_date = {r["due_date"]: r["n"] for r in cur.fetchall()}
    cur.execute("SELECT DATE(created_at) as d, COUNT(*) as n FROM events WHERE event_type='question_asked' GROUP BY DATE(created_at)")
    q_by_date = {r["d"]: r["n"] for r in cur.fetchall()}
    spikes = []
    for due_date, n_due in due_by_date.items():
        same_day = q_by_date.get(due_date, 0)
        prior_3 = sum(q_by_date.get(due_date - timedelta(days=k), 0) for k in range(1, 4))
        spikes.append({
            "due_date": due_date.isoformat(), "deadlines_due": n_due,
            "questions_same_day": same_day, "questions_prior_3_days": prior_3
        })
    spikes.sort(key=lambda x: x["due_date"])
    out["deadline_spikes"] = spikes[-30:]  # most recent/upcoming 30 dates with deadlines

    # ── Upload mix: permanent vs temporary vs admin/global ──
    cur.execute("""
        SELECT event_type, COUNT(*) as n FROM events
        WHERE event_type IN ('file_uploaded','temp_file_used','global_file_uploaded')
        GROUP BY event_type""")
    mix = {r["event_type"]: r["n"] for r in cur.fetchall()}
    out["upload_mix"] = {
        "permanent": mix.get("file_uploaded", 0),
        "temporary": mix.get("temp_file_used", 0),
        "global":    mix.get("global_file_uploaded", 0),
    }

    # ── Stale general-reference documents (untouched 90+ days) ──
    cur.execute("""
        SELECT id, orig_name, university, course as label,
               to_char(uploaded_at,'Mon DD YYYY') as uploaded_at
        FROM documents
        WHERE student_id IS NULL AND uploaded_at < NOW() - INTERVAL '90 days'
        ORDER BY uploaded_at ASC LIMIT 20""")
    out["stale_global_docs"] = [dict(r) for r in cur.fetchall()]

    # ── Rough "general reference material was available" rate: for each
    #    question, did that student's university have at least one global
    #    doc uploaded by that point? This is availability, not confirmed
    #    usage — we don't log whether the model's answer actually drew on it. ──
    cur.execute("""
        SELECT
          COUNT(*) FILTER (
            WHERE EXISTS (
              SELECT 1 FROM documents gd
              WHERE gd.student_id IS NULL
                AND lower(gd.university) = lower(st.university)
                AND gd.uploaded_at <= e.created_at
            )
          ) as with_docs,
          COUNT(*) as total
        FROM events e JOIN students st ON st.id = e.student_id
        WHERE e.event_type = 'question_asked'""")
    row = cur.fetchone()
    out["general_doc_availability_pct"] = (
        round(row["with_docs"] / row["total"] * 100, 1) if row and row["total"] else 0
    )

    return out

# ── Auth ──────────────────────────────────────────────────
@app.route("/")
def landing():
    # Always show landing page so students see the welcome screen first
    try:
        return render_template("landing.html")
    except Exception as e:
        print(f"landing error: {e}"); return render_template("landing.html")

@app.route("/register", methods=["GET","POST"])
def register():
    def err(msg):
        return render_template("register.html", error=msg,
                               classifications=CLASSIFICATIONS, majors=MAJORS)
    try:
        if request.method == "POST":
            email      = request.form.get("email","").strip().lower()
            pw         = request.form.get("password","").strip()
            fn         = request.form.get("first_name","").strip()
            ln         = request.form.get("last_name","").strip()
            cl         = request.form.get("classification","").strip()
            major      = request.form.get("major","").strip()
            university = request.form.get("university","").strip()
            if not all([email,pw,fn,ln,cl,major,university]):
                return err("All fields are required, including your university.")
            # WINK now serves students at any school, not just UTEP — the old
            # hardcoded @utep.edu/@miners.utep.edu check would lock every other
            # university's students out. A plain .edu check is a reasonable,
            # low-friction sanity check in its place; swap in a stricter
            # per-university domain check here later if you want to verify the
            # email actually matches the chosen school.
            if not email.endswith(".edu"):
                return err("Please use your school (.edu) email address.")
            if len(pw) < 8:
                return err("Password must be at least 8 characters.")
            if not DB_URL:
                return err("Database not configured.")
            conn = get_db(); cur = conn.cursor()
            cur.execute("SELECT id FROM students WHERE email=%s", (email,))
            if cur.fetchone():
                cur.close(); db_release(conn)
                return err("Account already exists — please log in.")
            cur.execute("""INSERT INTO students(email,password_hash,first_name,last_name,classification,major,university)
                           VALUES(%s,%s,%s,%s,%s,%s,%s) RETURNING id""",
                        (email, generate_password_hash(pw), fn, ln, cl, major, university))
            new_id = cur.fetchone()["id"]
            verify_token = secrets.token_urlsafe(32)
            cur.execute("UPDATE students SET verification_token=%s WHERE id=%s", (verify_token, new_id))
            conn.commit(); cur.close(); db_release(conn)
            session.permanent = True  # actually use the 7-day PERMANENT_SESSION_LIFETIME set above
            session["sid"] = new_id
            log_event(new_id, "account_created", {"email":email,"classification":cl,"major":major,"university":university})
            verify_link = url_for("verify_email", token=verify_token, _external=True)
            send_email(email, "Verify your WINK email address",
                       f"Hi {fn},\n\nWelcome to WINK! Please confirm your email address by visiting:\n{verify_link}\n\n"
                       f"You can use WINK right away either way — this just confirms we can reach you.\n\n— WINK")
            return redirect(url_for("documents"))
        return render_template("register.html", error=None,
                               classifications=CLASSIFICATIONS, majors=MAJORS)
    except Exception as e:
        print(f"register error: {e}"); traceback.print_exc()
        return err("Something went wrong on our end. Please try again in a moment.")

@app.route("/login", methods=["GET","POST"])
def login():
    try:
        if request.method == "POST":
            email = request.form.get("email","").strip().lower()
            pw    = request.form.get("password","").strip()
            if rate_limited(f"login:{request.remote_addr}", max_calls=10, window_seconds=60):
                return render_template("login.html", error="Too many attempts — please wait a minute and try again.")
            if not DB_URL:
                return render_template("login.html", error="Database not configured.")
            conn = get_db(); cur = conn.cursor()
            cur.execute("SELECT * FROM students WHERE email=%s", (email,))
            s = cur.fetchone(); cur.close(); db_release(conn)
            if s and check_password_hash(s["password_hash"], pw):
                if not s.get("is_active", True):
                    return render_template("login.html", error="This account has been suspended. Contact your administrator.")
                session.permanent = True  # actually use the 7-day PERMANENT_SESSION_LIFETIME set above
                session["sid"] = s["id"]
                log_event(s["id"], "login", {"email": email})
                # Admin goes straight to analytics
                if email == ADMIN_EMAIL:
                    return redirect(url_for("analytics_page"))
                return redirect(url_for("dashboard"))
            return render_template("login.html", error="Invalid email or password.")
        return render_template("login.html", error=None)
    except Exception as e:
        print(f"login error: {e}"); traceback.print_exc()
        return render_template("login.html", error="Something went wrong on our end. Please try again in a moment.")

@app.route("/logout")
def logout():
    session.clear(); return redirect(url_for("landing"))

@app.route("/verify-email/<token>")
def verify_email(token):
    """Confirms the email address behind a signup — doesn't gate access to
    WINK (the account works immediately at signup either way), it just marks
    email_verified so the admin dashboard can show which accounts have
    confirmed a reachable address."""
    if not DB_URL:
        return app.response_class("Database not configured.", mimetype="text/plain"), 500
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT id FROM students WHERE verification_token=%s", (token,))
        row = cur.fetchone()
        if not row:
            cur.close(); db_release(conn)
            return app.response_class(
                "This verification link is invalid or has already been used.",
                mimetype="text/plain"), 400
        cur.execute("UPDATE students SET email_verified=TRUE, verification_token=NULL WHERE id=%s", (row["id"],))
        conn.commit(); cur.close(); db_release(conn)
        log_event(row["id"], "email_verified", {})
        return app.response_class(
            "Your email is verified! You can close this tab and keep using WINK.",
            mimetype="text/plain")
    except Exception as e:
        print(f"verify_email error: {e}"); traceback.print_exc()
        return app.response_class("Something went wrong on our end. Please try again.", mimetype="text/plain"), 500

@app.route("/resend-verification", methods=["POST"])
def resend_verification():
    s = current_student()
    if not s: return jsonify({"error":"Not logged in"}), 401
    if s.get("email_verified"): return jsonify({"success": True, "already_verified": True})
    if rate_limited(f"resend-verify:{s['id']}", max_calls=3, window_seconds=300):
        return jsonify({"error": "Please wait a few minutes before requesting another email."}), 429
    try:
        token = secrets.token_urlsafe(32)
        conn = get_db(); cur = conn.cursor()
        cur.execute("UPDATE students SET verification_token=%s WHERE id=%s", (token, s["id"]))
        conn.commit(); cur.close(); db_release(conn)
        verify_link = url_for("verify_email", token=token, _external=True)
        send_email(s["email"], "Verify your WINK email address",
                   f"Hi {s['first_name']},\n\nPlease confirm your email address by visiting:\n{verify_link}\n\n— WINK")
        return jsonify({"success": True})
    except Exception as e:
        print(f"resend_verification error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500

@app.route("/forgot-password", methods=["POST"])
def forgot_password():
    email = request.form.get("email", "").strip().lower()
    try:
        if not email:
            return render_template("login.html", forgot=True, error="Please enter your email.")
        if rate_limited(f"forgot:{request.remote_addr}", max_calls=5, window_seconds=300):
            return render_template("login.html", forgot=True, error="Too many requests — please wait a few minutes and try again.")
        if not DB_URL:
            return render_template("login.html", forgot=True, error="Database not configured.")
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT id FROM students WHERE email=%s", (email,))
        s = cur.fetchone()
        if not s:
            cur.close(); db_release(conn)
            # Don't reveal whether the account exists — show the same confirmation either way.
            return render_template("login.html", forgot=True, reset_sent=True)

        token = secrets.token_urlsafe(32)
        expires_at = datetime.utcnow() + timedelta(hours=1)
        cur.execute(
            "INSERT INTO password_resets(student_id, token, expires_at) VALUES(%s,%s,%s)",
            (s["id"], token, expires_at)
        )
        conn.commit(); cur.close(); db_release(conn)
        log_event(s["id"], "password_reset_requested", {"email": email})

        reset_link = url_for("reset_password", token=token, _external=True)
        # Real email sending now — see send_email() / SMTP_* config above.
        # If no SMTP provider is configured, this falls back to logging the
        # link server-side (and only shows it in the response if
        # DEBUG_SHOW_RESET_LINKS is explicitly set for local testing) rather
        # than handing an account-takeover link to whoever submitted the form.
        sent = send_email(
            email, "Reset your WINK password",
            f"Hi,\n\nSomeone (hopefully you) requested a password reset for your WINK account.\n\n"
            f"Reset your password here: {reset_link}\n\nThis link expires in 1 hour. "
            f"If you didn't request this, you can safely ignore this email.\n\n— WINK"
        )
        if not sent:
            print(f"PASSWORD RESET LINK for {email}: {reset_link}")
        if not sent and DEBUG_SHOW_RESET_LINKS:
            return render_template("login.html", forgot=True, reset_sent=True, reset_link=reset_link)
        return render_template("login.html", forgot=True, reset_sent=True)
    except Exception as e:
        print(f"forgot_password error: {e}"); traceback.print_exc()
        return render_template("login.html", forgot=True, error="Something went wrong. Please try again.")

@app.route("/reset-password/<token>", methods=["GET", "POST"])
def reset_password(token):
    try:
        if not DB_URL:
            return render_template("login.html", error="Database not configured.")
        conn = get_db(); cur = conn.cursor()
        cur.execute("""
            SELECT pr.id AS reset_id, pr.student_id, pr.expires_at, pr.used, s.email
            FROM password_resets pr JOIN students s ON s.id = pr.student_id
            WHERE pr.token=%s
        """, (token,))
        row = cur.fetchone()

        if not row or row["used"] or row["expires_at"] < datetime.utcnow():
            cur.close(); db_release(conn)
            return render_template("reset_password.html", invalid=True)

        if request.method == "POST":
            pw = request.form.get("password", "").strip()
            confirm = request.form.get("confirm_password", "").strip()
            if len(pw) < 8:
                cur.close(); db_release(conn)
                return render_template("reset_password.html", token=token,
                                       error="Password must be at least 8 characters.")
            if pw != confirm:
                cur.close(); db_release(conn)
                return render_template("reset_password.html", token=token,
                                       error="Passwords do not match.")
            cur.execute("UPDATE students SET password_hash=%s WHERE id=%s",
                        (generate_password_hash(pw), row["student_id"]))
            cur.execute("UPDATE password_resets SET used=TRUE WHERE id=%s", (row["reset_id"],))
            conn.commit(); cur.close(); db_release(conn)
            log_event(row["student_id"], "password_reset_completed", {"email": row["email"]})
            return render_template("login.html", success="Your password has been updated — please sign in.")

        cur.close(); db_release(conn)
        return render_template("reset_password.html", token=token)
    except Exception as e:
        print(f"reset_password error: {e}"); traceback.print_exc()
        return render_template("reset_password.html", token=token, error="Something went wrong on our end. Please try again in a moment.")

# ── Pages ─────────────────────────────────────────────────
def get_questions_this_month(sid):
    if not DB_URL: return 0
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("""SELECT COUNT(*) as n FROM events
                       WHERE student_id=%s AND event_type='question_asked'
                       AND created_at >= date_trunc('month', NOW())""", (sid,))
        n = cur.fetchone()["n"]; cur.close(); db_release(conn)
        return n
    except Exception as e:
        print(f"get_questions_this_month error: {e}"); return 0

@app.route("/dashboard")
def dashboard():
    try:
        s = current_student()
        if not s: return redirect(url_for("login"))
        docs = get_docs(s["id"])
        upcoming_deadlines = get_upcoming_deadlines(s["id"], days_ahead=7)
        questions_this_month = get_questions_this_month(s["id"])
        log_event(s["id"], "page_view", {"page":"dashboard"})
        return render_template("dashboard.html", s=s, admin_email=ADMIN_EMAIL, docs=docs,
                               active="dashboard", max_docs=MAX_DOCS_PER_STUDENT,
                               upcoming_deadlines=upcoming_deadlines,
                               questions_this_month=questions_this_month)
    except Exception as e:
        print(f"dashboard error: {e}"); traceback.print_exc()
        return "<h2>Something went wrong</h2><p>Please try again, or <a href='/logout'>log out</a> and back in.</p>", 500

@app.route("/update-profile", methods=["POST"])
def update_profile():
    """Lets a student edit their own name, classification, and major from the
    dashboard. Email is intentionally not editable here — it's tied to login
    and to the ADMIN_EMAIL check elsewhere, so changing it needs more care
    than a quick profile edit."""
    try:
        s = current_student()
        if not s: return jsonify({"error": "Not logged in"}), 401
        data = request.get_json() or {}
        first_name     = (data.get("first_name") or "").strip()
        last_name      = (data.get("last_name") or "").strip()
        classification = (data.get("classification") or "").strip()
        major          = (data.get("major") or "").strip()
        university     = (data.get("university") or "").strip()
        if not all([first_name, last_name, classification, major, university]):
            return jsonify({"error": "All fields are required."}), 400
        if not DB_URL:
            return jsonify({"error": "No database configured."}), 500
        conn = get_db(); cur = conn.cursor()
        cur.execute("""UPDATE students SET first_name=%s, last_name=%s,
                       classification=%s, major=%s, university=%s WHERE id=%s""",
                    (first_name, last_name, classification, major, university, s["id"]))
        conn.commit(); cur.close(); db_release(conn)
        log_event(s["id"], "profile_updated", {"classification": classification, "major": major, "university": university})
        return jsonify({
            "success": True,
            "profile": {"first_name": first_name, "last_name": last_name,
                        "classification": classification, "major": major,
                        "university": university}
        })
    except Exception as e:
        print(f"update_profile error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500

@app.route("/documents")
def documents():
    try:
        s = current_student()
        if not s: return redirect(url_for("login"))
        docs = get_docs(s["id"])
        grouped_docs = group_docs_by_course(docs)
        known_courses = sorted({(d.get("course") or "").strip() for d in docs
                                 if (d.get("course") or "").strip()})
        log_event(s["id"], "page_view", {"page":"documents"})
        return render_template("documents.html", s=s, admin_email=ADMIN_EMAIL, docs=docs,
                               grouped_docs=grouped_docs, known_courses=known_courses,
                               active="documents", max_docs=MAX_DOCS_PER_STUDENT)
    except Exception as e:
        print(f"documents error: {e}"); traceback.print_exc()
        return "<h2>Something went wrong</h2><p>Please try again, or <a href='/logout'>log out</a> and back in.</p>", 500

@app.route("/chat-page")
def chat_page():
    try:
        s = current_student()
        if not s: return redirect(url_for("login"))
        docs = get_docs(s["id"])
        log_event(s["id"], "page_view", {"page":"chat"})
        return render_template("chat.html", s=s, admin_email=ADMIN_EMAIL, docs=docs, active="chat")
    except Exception as e:
        print(f"chat_page error: {e}"); traceback.print_exc()
        return "<h2>Something went wrong</h2><p>Please try again, or <a href='/logout'>log out</a> and back in.</p>", 500

@app.route("/analytics-page")
def analytics_page():
    try:
        s = current_student()
        if not s: return redirect(url_for("login"))
        if s["email"].lower() != ADMIN_EMAIL:
            return redirect(url_for("dashboard"))
        log_event(s["id"], "page_view", {"page":"analytics"})
        return render_template("analytics.html", s=s, admin_email=ADMIN_EMAIL, active="analytics")
    except Exception as e:
        print(f"analytics_page error: {e}"); traceback.print_exc()
        return "<h2>Something went wrong</h2><p>Please try again, or <a href='/logout'>log out</a> and back in.</p>", 500

# ── API ───────────────────────────────────────────────────
@app.route("/upload", methods=["POST"])
def upload_file():
    try:
        s = current_student()
        if not s: return jsonify({"error":"Not logged in"}), 401
        wait = rate_limited(f"upload:{s['id']}", max_calls=10, window_seconds=60)
        if wait:
            return jsonify({"error": "Too many uploads in a row — please wait a moment.", "retry_after": wait}), 429
        if "file" not in request.files:
            return jsonify({"error":"No file"}), 400
        file      = request.files["file"]
        temporary = request.form.get("temporary","").strip().lower() == "true"
        if not file or not file.filename:
            return jsonify({"error":"No file selected"}), 400
        ext = file.filename.rsplit(".",1)[-1].lower() if "." in file.filename else ""
        if ext not in ALLOWED_EXT:
            return jsonify({"error":f"File type .{ext} not allowed"}), 400

        # Temporary, this-conversation-only upload: extract the text and hand
        # it straight back to the client — never written to the documents
        # table, so it doesn't count against MAX_DOCS_PER_STUDENT and never
        # shows up in My Documents. The client resends this content with
        # each /chat call for the current conversation only; nothing here
        # persists once that conversation ends.
        if temporary:
            tmp_name = f"{uuid.uuid4().hex[:8]}_{secure_filename(file.filename)}"
            tmp_path = os.path.join(UPLOAD_FOLDER, tmp_name)
            try:
                file.save(tmp_path)
                content = extract_text(tmp_path, file.filename)
            finally:
                if os.path.exists(tmp_path):
                    try: os.remove(tmp_path)
                    except Exception: pass
            content = content[:MAX_TEMP_DOC_CHARS]
            log_event(s["id"], "temp_file_used", {"name": file.filename, "chars": len(content)})
            return jsonify({
                "success": True, "temporary": True,
                "name": file.filename, "content": content,
                "chars_extracted": len(content)
            })

        course = request.form.get("course","").strip()
        crn    = request.form.get("crn","").strip()
        if not course:
            return jsonify({"error":"Please enter a course name."}), 400
        if not crn:
            return jsonify({"error":"Please enter a CRN#."}), 400

        replaced = False
        if DB_URL:
            # Document versioning: re-uploading the same filename for the same
            # course + CRN replaces the old copy instead of adding a new one —
            # this is almost always "the professor updated the syllabus," not
            # "a 21st document," and it keeps students from hitting the cap
            # just from re-uploading a corrected file.
            conn = get_db(); cur = conn.cursor()
            cur.execute("""SELECT id, filename FROM documents
                           WHERE student_id=%s AND lower(course)=lower(%s)
                           AND crn=%s AND lower(orig_name)=lower(%s)""",
                        (s["id"], course, crn, file.filename))
            existing = cur.fetchone()
            if existing:
                old_fp = os.path.join(UPLOAD_FOLDER, str(s["id"]), existing["filename"])
                if os.path.exists(old_fp):
                    try: os.remove(old_fp)
                    except Exception: pass
                cur.execute("DELETE FROM documents WHERE id=%s", (existing["id"],))
                conn.commit()
                replaced = True
            cur.close(); db_release(conn)

        if not replaced:
            existing_docs = get_docs(s["id"])
            if len(existing_docs) >= MAX_DOCS_PER_STUDENT:
                return jsonify({
                    "error": f"You've reached the {MAX_DOCS_PER_STUDENT}-document limit. "
                             f"Delete a document before uploading a new one."
                }), 400

        folder = os.path.join(UPLOAD_FOLDER, str(s["id"]))
        os.makedirs(folder, exist_ok=True)
        orig  = file.filename
        saved = f"{uuid.uuid4().hex[:8]}_{secure_filename(orig)}"
        path  = os.path.join(folder, saved)
        file.save(path)
        size    = os.path.getsize(path)
        content = extract_text(path, orig)
        print(f"UPLOAD: {orig} → {len(content)} chars extracted")
        new_doc_id = None
        if DB_URL:
            conn = get_db(); cur = conn.cursor()
            cur.execute("""INSERT INTO documents
                           (student_id,filename,orig_name,course,crn,size_bytes,content)
                           VALUES(%s,%s,%s,%s,%s,%s,%s) RETURNING id""",
                        (s["id"], saved, orig, course, crn, size, content))
            new_doc_id = cur.fetchone()["id"]
            conn.commit(); cur.close(); db_release(conn)

        # Deadline extraction: one small Haiku call per upload to pull out
        # assignment/exam dates so they can show up on the dashboard and in
        # reminder emails. Best-effort — never blocks the upload if it fails.
        deadlines_found = 0
        if new_doc_id and content:
            deadlines = extract_deadlines(content)
            if deadlines and DB_URL:
                conn = get_db(); cur = conn.cursor()
                for d in deadlines:
                    cur.execute("""INSERT INTO deadlines(student_id,document_id,course,title,due_date)
                                   VALUES(%s,%s,%s,%s,%s)""",
                                (s["id"], new_doc_id, course, d["title"], d["due_date"]))
                conn.commit(); cur.close(); db_release(conn)
                deadlines_found = len(deadlines)

        log_event(s["id"], "file_replaced" if replaced else "file_uploaded",
                  {"name":orig,"course":course,"crn":crn,"chars":len(content),"deadlines":deadlines_found})
        return jsonify({
            "success":True, "docs":get_docs(s["id"]), "chars_extracted":len(content),
            "replaced": replaced, "deadlines_found": deadlines_found
        })
    except Exception as e:
        print(f"upload error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500

@app.route("/delete-file", methods=["POST"])
def delete_file():
    try:
        s = current_student()
        if not s: return jsonify({"error":"Not logged in"}), 401
        doc_id = (request.get_json() or {}).get("doc_id")
        if DB_URL and doc_id:
            conn = get_db(); cur = conn.cursor()
            cur.execute("SELECT filename FROM documents WHERE id=%s AND student_id=%s", (doc_id, s["id"]))
            doc = cur.fetchone()
            if doc:
                fp = os.path.join(UPLOAD_FOLDER, str(s["id"]), doc["filename"])
                if os.path.exists(fp): os.remove(fp)
                cur.execute("DELETE FROM documents WHERE id=%s", (doc_id,))
                conn.commit()
                log_event(s["id"], "file_deleted", {"doc_id": doc_id})
            cur.close(); db_release(conn)
        return jsonify({"success":True, "docs":get_docs(s["id"])})
    except Exception as e:
        print(f"delete error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500

# ── General reference documents (admin-only) ────────────────
# These apply to every student's chat automatically (see build_global_doc_context
# and its use in /chat) but are stored with student_id=NULL, so they never show
# up in any student's own "My Documents" list or count against their 20-doc cap.
@app.route("/global-documents")
def list_global_documents():
    s = current_student()
    if not s or s["email"].lower() != ADMIN_EMAIL: return jsonify({"error":"Not authorized"}), 403
    university = request.args.get("university","").strip()
    return jsonify({"docs": get_global_docs(university or None)})

@app.route("/upload-global", methods=["POST"])
def upload_global_document():
    try:
        s = current_student()
        if not s or s["email"].lower() != ADMIN_EMAIL:
            return jsonify({"error":"Not authorized"}), 403
        if "file" not in request.files:
            return jsonify({"error":"No file"}), 400
        file       = request.files["file"]
        label      = request.form.get("label","").strip() or "General"
        university = request.form.get("university","").strip()
        if not university:
            return jsonify({"error":"Please choose which university this document applies to."}), 400
        if not file or not file.filename:
            return jsonify({"error":"No file selected"}), 400
        ext = file.filename.rsplit(".",1)[-1].lower() if "." in file.filename else ""
        if ext not in ALLOWED_EXT:
            return jsonify({"error":f"File type .{ext} not allowed"}), 400

        folder = os.path.join(UPLOAD_FOLDER, "global")
        os.makedirs(folder, exist_ok=True)
        orig  = file.filename
        saved = f"{uuid.uuid4().hex[:8]}_{secure_filename(orig)}"
        path  = os.path.join(folder, saved)
        file.save(path)
        size    = os.path.getsize(path)
        content = extract_text(path, orig)
        print(f"GLOBAL UPLOAD: {orig} ({university}) → {len(content)} chars extracted")
        if DB_URL:
            conn = get_db(); cur = conn.cursor()
            cur.execute("""INSERT INTO documents
                           (student_id,filename,orig_name,course,crn,size_bytes,content,university)
                           VALUES(NULL,%s,%s,%s,'',%s,%s,%s)""",
                        (saved, orig, label, size, content, university))
            conn.commit(); cur.close(); db_release(conn)
        log_event(s["id"], "global_file_uploaded", {"name": orig, "label": label, "university": university, "chars": len(content)})
        return jsonify({"success": True, "docs": get_global_docs(university), "chars_extracted": len(content)})
    except Exception as e:
        print(f"global upload error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500

@app.route("/delete-global-document", methods=["POST"])
def delete_global_document():
    try:
        s = current_student()
        if not s or s["email"].lower() != ADMIN_EMAIL: return jsonify({"error":"Not authorized"}), 403
        data       = request.get_json() or {}
        doc_id     = data.get("doc_id")
        university = (data.get("university") or "").strip()
        if DB_URL and doc_id:
            conn = get_db(); cur = conn.cursor()
            cur.execute("SELECT filename, university FROM documents WHERE id=%s AND student_id IS NULL", (doc_id,))
            doc = cur.fetchone()
            if doc:
                fp = os.path.join(UPLOAD_FOLDER, "global", doc["filename"])
                if os.path.exists(fp): os.remove(fp)
                cur.execute("DELETE FROM documents WHERE id=%s", (doc_id,))
                conn.commit()
                log_event(s["id"], "global_file_deleted", {"doc_id": doc_id})
            cur.close(); db_release(conn)
        return jsonify({"success": True, "docs": get_global_docs(university or None)})
    except Exception as e:
        print(f"delete global error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500

def get_upcoming_deadlines(sid, days_ahead=14):
    if not DB_URL: return []
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("""SELECT id, course, title, due_date FROM deadlines
                       WHERE student_id=%s AND due_date >= CURRENT_DATE
                       AND due_date <= CURRENT_DATE + %s::int
                       ORDER BY due_date ASC""", (sid, days_ahead))
        rows = [dict(r) for r in cur.fetchall()]
        cur.close(); db_release(conn)
        for r in rows:
            r["due_date"] = r["due_date"].isoformat()
        return rows
    except Exception as e:
        print(f"get_upcoming_deadlines error: {e}"); return []

def get_all_deadlines(sid):
    """Every deadline row for a student, no date-range cap — the single
    source of truth used by both the chat context (as text) and the visual
    calendar page (as JSON)."""
    if not DB_URL: return []
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("""SELECT dl.id, dl.course, dl.title, dl.due_date,
                              dl.document_id, d.orig_name as document_name
                       FROM deadlines dl LEFT JOIN documents d ON d.id = dl.document_id
                       WHERE dl.student_id=%s ORDER BY dl.due_date ASC""", (sid,))
        rows = [dict(r) for r in cur.fetchall()]
        cur.close(); db_release(conn)
        for r in rows:
            r["due_date"] = r["due_date"].isoformat() if r["due_date"] else None
        return rows
    except Exception as e:
        print(f"get_all_deadlines error: {e}"); return []

def build_deadlines_context(sid):
    """Every deadline extracted from every one of the student's uploaded
    documents, across every course, with no date-range cap and no truncation.
    This is deliberately separate from build_doc_context()'s truncated raw
    document text — a "build me a master calendar" question shouldn't depend
    on how much raw document text fit under the per-message cost cap, since
    the structured deadline data is already small and already complete."""
    if not DB_URL:
        return "\n\nNo deadline data available (no database configured)."
    rows = get_all_deadlines(sid)
    if not rows:
        return ("\n\nNo deadlines have been extracted yet. This can mean the student's "
                "documents don't contain a schedule of specific dates, or nothing has "
                "been uploaded yet — don't invent dates that aren't in this list.")
    lines = [f"\n\n{'='*60}\nEXTRACTED DEADLINES — every date-specific item found across "
             f"ALL of the student's uploaded documents ({len(rows)} total). This list is "
             "COMPLETE and NOT truncated, unlike the raw document text below — always use "
             "this list (not the raw text) when asked for a calendar, schedule, or 'what's "
             f"due' summary.\n{'='*60}"]
    for r in rows:
        due = r["due_date"] or "date unknown"
        lines.append(f"- [{r['course']}] {r['title']} — due {due}")
    lines.append(f"{'='*60}\n")
    return "\n".join(lines)

@app.route("/deadlines")
def deadlines():
    s = current_student()
    if not s: return jsonify({"error":"Not logged in"}), 401
    days = min(int(request.args.get("days", 14)), 90)
    return jsonify({"deadlines": get_upcoming_deadlines(s["id"], days)})

@app.route("/calendar-page")
def calendar_page():
    s = current_student()
    if not s: return redirect(url_for("login"))
    log_event(s["id"], "page_view", {"page":"calendar"})
    return render_template("calendar.html", s=s, admin_email=ADMIN_EMAIL, active="calendar")

@app.route("/calendar-data")
def calendar_data():
    s = current_student()
    if not s: return jsonify({"error":"Not logged in"}), 401
    return jsonify({"deadlines": get_all_deadlines(s["id"])})

@app.route("/reprocess-deadlines", methods=["POST"])
def reprocess_deadlines():
    """Re-run deadline extraction against documents that are already stored
    (using their already-extracted, already-stored text — no re-upload
    needed). Useful for documents uploaded before DEADLINE_EXTRACTION_MAX_CHARS
    was fixed, or if a document's schedule wasn't picked up the first time."""
    s = current_student()
    if not s: return jsonify({"error":"Not logged in"}), 401
    if not DB_URL: return jsonify({"error":"No database"}), 500
    if rate_limited(f"reprocess:{s['id']}", max_calls=3, window_seconds=300):
        return jsonify({"error": "Please wait a few minutes before doing this again."}), 429
    try:
        docs = get_docs(s["id"])
        total_found = 0
        docs_processed = 0
        for d in docs:
            content = (d.get("content") or "").strip()
            if not content:
                continue
            found = extract_deadlines(content)
            conn = get_db(); cur = conn.cursor()
            # Replace this document's deadlines rather than duplicating them
            cur.execute("DELETE FROM deadlines WHERE document_id=%s", (d["id"],))
            for item in found:
                cur.execute("""INSERT INTO deadlines(student_id,document_id,course,title,due_date)
                               VALUES(%s,%s,%s,%s,%s)""",
                            (s["id"], d["id"], d["course"], item["title"], item["due_date"]))
            conn.commit(); cur.close(); db_release(conn)
            docs_processed += 1
            total_found += len(found)
        log_event(s["id"], "deadlines_reprocessed", {"docs": docs_processed, "found": total_found})
        return jsonify({"success": True, "documents_processed": docs_processed, "deadlines_found": total_found})
    except Exception as e:
        print(f"reprocess_deadlines error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end."}), 500

@app.route("/send-deadline-reminders", methods=["POST"])
def send_deadline_reminders():
    """Meant to be hit once a day by an external scheduler (Render cron job,
    GitHub Action, etc.) with ?key=CRON_SECRET — emails each student a
    summary of anything due in the next 3 days that they haven't already
    been reminded about."""
    if not CRON_SECRET or request.args.get("key") != CRON_SECRET:
        return jsonify({"error": "Not authorized"}), 403
    if not DB_URL:
        return jsonify({"error": "No database"}), 500
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("""SELECT d.id, d.title, d.due_date, d.course, s.id as sid, s.email, s.first_name
                       FROM deadlines d JOIN students s ON s.id = d.student_id
                       WHERE d.reminded = FALSE
                       AND d.due_date BETWEEN CURRENT_DATE AND CURRENT_DATE + 3
                       ORDER BY s.id, d.due_date""")
        rows = [dict(r) for r in cur.fetchall()]
        cur.close(); db_release(conn)

        by_student = {}
        for r in rows:
            by_student.setdefault(r["sid"], {"email": r["email"], "first_name": r["first_name"], "items": []})
            by_student[r["sid"]]["items"].append(r)

        sent_count = 0
        for sid, info in by_student.items():
            lines = [f"  • {it['title']} ({it['course']}) — due {it['due_date'].strftime('%A, %b %d')}"
                     for it in info["items"]]
            body = (f"Hi {info['first_name']},\n\nHere's what's coming up in the next few days:\n\n"
                    + "\n".join(lines) + "\n\n— WINK")
            if send_email(info["email"], "Upcoming deadlines — WINK", body):
                sent_count += 1

        if rows:
            ids = [r["id"] for r in rows]
            conn = get_db(); cur = conn.cursor()
            cur.execute("UPDATE deadlines SET reminded=TRUE WHERE id = ANY(%s)", (ids,))
            conn.commit(); cur.close(); db_release(conn)

        return jsonify({"students_notified": sent_count, "deadlines_covered": len(rows)})
    except Exception as e:
        print(f"send_deadline_reminders error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end."}), 500

@app.route("/chat", methods=["POST"])
def chat():
    try:
        s = current_student()
        if not s: return jsonify({"error":"Not logged in"}), 401
        if not ANTHROPIC_API_KEY:
            return jsonify({"error":"ANTHROPIC_API_KEY not set"}), 500
        wait = rate_limited(f"chat:{s['id']}", max_calls=20, window_seconds=60)
        if wait:
            return jsonify({
                "error": "You're asking questions faster than I can keep up — please wait a moment and try again.",
                "retry_after": wait
            }), 429
        data     = request.get_json() or {}
        messages = data.get("messages", [])
        user_msg = messages[-1]["content"] if messages else ""
        if isinstance(user_msg, str) and len(user_msg) > MAX_USER_MESSAGE_CHARS:
            return jsonify({"error": f"That message is too long (max {MAX_USER_MESSAGE_CHARS} characters). Please shorten it and try again."}), 400
        log_event(s["id"], "question_asked", {"q": user_msg[:200]})

        # Conversation history (feature): load or create the conversation this
        # message belongs to, so it can be revisited/exported later. A client
        # that doesn't send conversation_id still works exactly as before —
        # this just also saves a copy server-side.
        conv_id = data.get("conversation_id")
        conv_row = None
        if DB_URL:
            conn = get_db(); cur = conn.cursor()
            if conv_id:
                cur.execute("SELECT id, title, messages FROM conversations WHERE id=%s AND student_id=%s",
                            (conv_id, s["id"]))
                conv_row = cur.fetchone()
            if not conv_row:
                title = (str(user_msg).strip()[:60] or "New conversation")
                cur.execute("""INSERT INTO conversations(student_id, title, messages)
                               VALUES(%s,%s,'[]') RETURNING id, title, messages""", (s["id"], title))
                conv_row = cur.fetchone()
                conn.commit()
            cur.close(); db_release(conn)
            conv_id = conv_row["id"]

        # Cost control: don't let conversation history grow unbounded —
        # everything sent here gets re-billed as input tokens every turn.
        messages = messages[-MAX_CHAT_HISTORY_MESSAGES:]
        while messages and messages[0].get("role") != "user":
            messages.pop(0)

        docs         = get_docs(s["id"])
        doc_ctx      = build_doc_context(docs)
        deadline_ctx = build_deadlines_context(s["id"])
        student_university = (s.get("university") or "").strip()
        global_ctx   = build_global_doc_context(get_global_docs(student_university or None), student_university)

        # Temporary, this-conversation-only file (see /upload's `temporary`
        # flag). The client resends the extracted content with every message
        # in this conversation — nothing here is read from or written to the
        # documents table, so it's never saved and never counts toward the
        # student's upload cap.
        temp_doc = data.get("temp_doc")
        if isinstance(temp_doc, dict) and temp_doc.get("content"):
            t_name = str(temp_doc.get("name") or "attached file")[:200]
            t_content = str(temp_doc["content"])[:MAX_TEMP_DOC_CHARS]
            temp_doc_ctx = (
                f"\n\nThe student has temporarily attached a file for THIS CONVERSATION "
                f"ONLY (not saved to their account, not one of their uploaded documents): "
                f"'{t_name}'.\n\n{t_content}"
            )
        else:
            temp_doc_ctx = ""
        import datetime
        now   = datetime.datetime.now()
        today = now.strftime("%A, %B %d, %Y")
        university_display = student_university or "their university"
        is_utep = "utep" in student_university.lower() or "el paso" in student_university.lower()
        instructions = (
            f"You are WINK, a warm encouraging AI-powered Academic Support System for college students. "
            f"Today's date is {today}. Always use this when answering questions about "
            f"deadlines, schedules, or anything time-related. "
            f"You are helping {s['first_name']} {s['last_name']}, "
            f"a {s['classification']} majoring in {s['major']} at {university_display}. "
            f"ANSWERING STRATEGY — follow this order: "
            f"1. For calendars, schedules, or 'what's due' questions across any or all "
            f"courses, use the EXTRACTED DEADLINES list below — it is complete and not "
            f"truncated. Do not tell the student their documents were truncated when "
            f"answering these — the deadlines list already accounts for that. "
            f"2. For anything else, check the student's uploaded documents below for the "
            f"answer — every uploaded document appears there, so never tell the student to "
            f"re-upload something that's already listed. If found, quote directly from their "
            f"documents with specific details. "
            + (f"2b. Also check the file the student temporarily attached to THIS "
               f"conversation below — answer questions about it the same way you would "
               f"an uploaded document. If they ask about saving it for later, let them "
               f"know it's only available in this conversation, and they can upload it "
               f"permanently instead from the My Documents page if they want to keep it. "
               if temp_doc_ctx else "")
            + f"3. Also check the GENERAL REFERENCE DOCUMENTS block below, which applies to every "
            f"student — use it the same way, but don't call it 'your document' since the "
            f"student didn't upload it themselves. "
            f"4. If the answer is NOT in either of those, use the web_search tool "
            f"to find current, accurate information from the internet — always search "
            f"specifically for {university_display} when the question is campus-specific "
            f"(e.g. \"{university_display} writing center hours\", not just \"writing center "
            f"hours\"). "
            f"This includes questions about professors, university staff, campus resources, "
            f"current events, university policies, people at the university, and anything "
            f"not covered in their uploaded files. "
            f"5. Always tell the student whether your answer came from their documents, "
            f"the extracted deadlines list, or a web search, so they know the source — but "
            f"never mention the GENERAL REFERENCE DOCUMENTS as a separate source out loud. "
            f"RICH CONTENT — the chat interface CAN render maps and images, so use them "
            f"whenever they'd genuinely help: "
            f"- For a campus building, address, or any physical location, include "
            f"[[map: specific place name or address]] on its own — e.g. [[map: Union "
            f"Building, {university_display}]]. Always add the university name (and its "
            f"city/state if you know it) to the query so the map centers on the right place. "
            f"Do not say you can't show a map — you can, using this syntax. "
            f"- For a photo of a real, notable person, place, or subject that likely has a "
            f"Wikipedia page (e.g. a university president, a historical figure, a well-known "
            f"landmark), include [[image: subject name]] on its own — e.g. [[image: Heather "
            f"Wilson]]. This looks up a real photo automatically; don't say you can't show a "
            f"picture, use this instead. It will show 'no photo found' on its own if the "
            f"subject doesn't have one — you don't need to hedge about that in your text. "
            f"This only works for subjects with a public Wikipedia page — it will not find "
            f"photos of specific campus buildings, ordinary people, or anything from the "
            f"student's own documents; don't use it for those. "
            f"- For any other image, use standard markdown ![description](image URL) — only "
            f"with a real URL you found via web_search, never a URL you're guessing at or "
            f"making up. If you don't have a real image URL from search results, don't "
            f"fabricate one — just answer with text instead. "
            + ("UTEP president is Heather Wilson. UTEP resources: University Writing Center, "
               "CASS Tutoring, Advising & Student Support. "
               if is_utep else
               f"For {university_display}-specific facts (current president, named campus "
               f"resources, offices, etc.), use web_search rather than guessing — don't assume "
               f"UTEP's resources or leadership apply here. ")
            + f"TONE: Be warm, specific, actionable, and confident. Never narrate your own "
            f"process out loud — don't say things like 'I'll look for that' or 'let me try to "
            f"find that' or 'I'll search for it'; just do it and answer with what you found, "
            f"stated plainly as fact. Use at most 2-3 emoji per answer, placed where they "
            f"genuinely add warmth or clarity (e.g. next to a heading, an encouraging line, or "
            f"a key point) — never more than that, and never one on every line. "
            f"CRITICAL THINKING & GROWTH MINDSET: don't just hand over the answer and stop. "
            f"Where it fits naturally, add a short follow-up that pushes the student's thinking "
            f"further — e.g. ask them to predict the next step before you confirm it, suggest "
            f"they explain the concept back in their own words, point out a related question "
            f"worth exploring, connect the topic to something they already know, or note what "
            f"they should try themselves before asking again next time. Keep this brief (one "
            f"sentence is usually enough) and vary it — don't repeat the same prompt every "
            f"time or force it into an answer where it doesn't fit. End with an encouraging note."
        )
        # Cost control: mark the (large, per-student-static) document context as
        # cacheable. Anthropic bills cache reads at roughly 10% of the standard
        # input rate, so any second-or-later question in the same session costs
        # far less on this block instead of re-billing it at full price every time.
        # Each context block gets its own cache breakpoint — updating one (e.g.
        # a new upload) doesn't invalidate the cache on the others.
        system = [
            {"type": "text", "text": instructions},
            {"type": "text", "text": deadline_ctx, "cache_control": {"type": "ephemeral"}},
            {"type": "text", "text": global_ctx, "cache_control": {"type": "ephemeral"}},
            {"type": "text", "text": doc_ctx, "cache_control": {"type": "ephemeral"}},
        ]
        if temp_doc_ctx:
            system.append({"type": "text", "text": temp_doc_ctx, "cache_control": {"type": "ephemeral"}})
        # Speed: reuse the single client created at module load (see top of
        # file) instead of opening a brand-new connection for every question.
        client = anthropic_client
        if client is None:
            return jsonify({"error":"ANTHROPIC_API_KEY not set"}), 500

        def generate():
            full_reply = []
            try:
                with client.messages.stream(
                    model=CHAT_MODEL,
                    max_tokens=CHAT_MAX_TOKENS,
                    system=system,
                    messages=messages,
                    tools=[{"type": "web_search_20250305", "name": "web_search", "max_uses": WEB_SEARCH_MAX_USES}]
                ) as stream:
                    for text in stream.text_stream:
                        full_reply.append(text)
                        yield text
            except Exception as e:
                print(f"stream error: {e}"); traceback.print_exc()
                yield "\n\nSomething went wrong on our end — please try asking again."
            reply = "".join(full_reply) or "I had trouble finding an answer — please try again."
            log_event(s["id"], "answer_given", {"len": len(reply), "full_answer": reply})
            if DB_URL and conv_id:
                try:
                    conn = get_db(); cur = conn.cursor()
                    saved = safe_payload(conv_row["messages"]) if isinstance(conv_row["messages"], str) else (conv_row["messages"] or [])
                    if not isinstance(saved, list): saved = []
                    saved.append({"role": "user", "content": user_msg, "ts": datetime.datetime.utcnow().isoformat()})
                    saved.append({"role": "assistant", "content": reply, "ts": datetime.datetime.utcnow().isoformat()})
                    cur.execute("UPDATE conversations SET messages=%s, updated_at=NOW() WHERE id=%s",
                                (json.dumps(saved), conv_id))
                    conn.commit(); cur.close(); db_release(conn)
                except Exception as e:
                    print(f"conversation save error: {e}")

        resp = app.response_class(generate(), mimetype="text/plain")
        resp.headers["X-Accel-Buffering"] = "no"
        resp.headers["Cache-Control"] = "no-cache"
        if conv_id:
            resp.headers["X-Conversation-Id"] = str(conv_id)
        return resp
    except Exception as e:
        print(f"chat error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500


@app.route("/conversations")
def list_conversations():
    s = current_student()
    if not s: return jsonify({"error":"Not logged in"}), 401
    if not DB_URL: return jsonify({"conversations": []})
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("""SELECT id, title, messages, updated_at FROM conversations
                       WHERE student_id=%s ORDER BY updated_at DESC LIMIT 50""", (s["id"],))
        rows = cur.fetchall(); cur.close(); db_release(conn)
        out = []
        for r in rows:
            msgs = safe_payload(r["messages"]) if isinstance(r["messages"], str) else (r["messages"] or [])
            out.append({
                "id": r["id"], "title": r["title"],
                "updated_at": r["updated_at"].isoformat() if r["updated_at"] else None,
                "message_count": len(msgs) if isinstance(msgs, list) else 0,
            })
        return jsonify({"conversations": out})
    except Exception as e:
        print(f"list_conversations error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end."}), 500

@app.route("/conversations/<int:conv_id>")
def get_conversation(conv_id):
    s = current_student()
    if not s: return jsonify({"error":"Not logged in"}), 401
    if not DB_URL: return jsonify({"error":"No database"}), 500
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT id, title, messages FROM conversations WHERE id=%s AND student_id=%s",
                    (conv_id, s["id"]))
        row = cur.fetchone(); cur.close(); db_release(conn)
        if not row: return jsonify({"error": "Not found"}), 404
        msgs = safe_payload(row["messages"]) if isinstance(row["messages"], str) else (row["messages"] or [])
        return jsonify({"id": row["id"], "title": row["title"], "messages": msgs})
    except Exception as e:
        print(f"get_conversation error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end."}), 500

@app.route("/conversations/<int:conv_id>/delete", methods=["POST"])
def delete_conversation(conv_id):
    s = current_student()
    if not s: return jsonify({"error":"Not logged in"}), 401
    if not DB_URL: return jsonify({"error":"No database"}), 500
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("DELETE FROM conversations WHERE id=%s AND student_id=%s", (conv_id, s["id"]))
        conn.commit(); cur.close(); db_release(conn)
        return jsonify({"success": True})
    except Exception as e:
        print(f"delete_conversation error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end."}), 500

def _conversation_transcript(title, msgs):
    lines = [f"# {title}", ""]
    for m in msgs:
        who = "You" if m.get("role") == "user" else "WINK"
        lines.append(f"**{who}:** {m.get('content','')}")
        lines.append("")
    return "\n".join(lines)

@app.route("/conversations/<int:conv_id>/export")
def export_conversation(conv_id):
    s = current_student()
    if not s: return jsonify({"error":"Not logged in"}), 401
    if not DB_URL: return jsonify({"error":"No database"}), 500
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT title, messages FROM conversations WHERE id=%s AND student_id=%s",
                    (conv_id, s["id"]))
        row = cur.fetchone(); cur.close(); db_release(conn)
        if not row: return jsonify({"error": "Not found"}), 404
        msgs = safe_payload(row["messages"]) if isinstance(row["messages"], str) else (row["messages"] or [])
        transcript = _conversation_transcript(row["title"], msgs)
        resp = app.response_class(transcript, mimetype="text/markdown")
        safe_name = secure_filename(row["title"])[:40] or "conversation"
        resp.headers["Content-Disposition"] = f'attachment; filename="{safe_name}.md"'
        return resp
    except Exception as e:
        print(f"export_conversation error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end."}), 500

@app.route("/conversations/<int:conv_id>/share", methods=["POST"])
def share_conversation(conv_id):
    s = current_student()
    if not s: return jsonify({"error":"Not logged in"}), 401
    if not DB_URL: return jsonify({"error":"No database"}), 500
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT id, share_token FROM conversations WHERE id=%s AND student_id=%s",
                    (conv_id, s["id"]))
        row = cur.fetchone()
        if not row:
            cur.close(); db_release(conn)
            return jsonify({"error": "Not found"}), 404
        token = row["share_token"] or secrets.token_urlsafe(24)
        if not row["share_token"]:
            cur.execute("UPDATE conversations SET share_token=%s WHERE id=%s", (token, conv_id))
            conn.commit()
        cur.close(); db_release(conn)
        return jsonify({"share_url": url_for("view_shared_conversation", token=token, _external=True)})
    except Exception as e:
        print(f"share_conversation error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end."}), 500

@app.route("/shared/<token>")
def view_shared_conversation(token):
    """Public, read-only view of a shared conversation. No login required —
    the unguessable token is the access control, same model as a shared
    Google Doc link."""
    if not DB_URL: return "Not available.", 404
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT title, messages FROM conversations WHERE share_token=%s", (token,))
        row = cur.fetchone(); cur.close(); db_release(conn)
        if not row: return "This shared conversation could not be found.", 404
        msgs = safe_payload(row["messages"]) if isinstance(row["messages"], str) else (row["messages"] or [])
        rows_html = "".join(
            f'<div style="margin-bottom:14px;"><strong style="color:{"#FF8200" if m.get("role")=="user" else "#002855"};">'
            f'{"You" if m.get("role")=="user" else "WINK"}:</strong> '
            f'<span style="white-space:pre-wrap;">{(m.get("content","") or "").replace("<","&lt;").replace(">","&gt;")}</span></div>'
            for m in msgs
        )
        return f"""<!DOCTYPE html><html><head><meta charset="UTF-8">
<title>{row['title']} — WINK</title></head>
<body style="font-family:-apple-system,sans-serif;max-width:700px;margin:40px auto;padding:0 20px;color:#444;">
<h1 style="color:#002855;">{row['title']}</h1>
<p style="color:#6b7a99;font-size:13px;">Shared read-only conversation from WINK</p>
<hr style="border:none;border-top:1px solid #dde3f0;margin:20px 0;">
{rows_html}
</body></html>"""
    except Exception as e:
        print(f"view_shared_conversation error: {e}"); traceback.print_exc()
        return "Something went wrong.", 500

@app.route("/analytics-data")
def analytics_data():
    try:
        s = current_student()
        if not s: return jsonify({"error":"Not logged in"}), 401
        if s["email"].lower() != ADMIN_EMAIL: return jsonify({"error":"Not authorized"}), 403
        if not DB_URL: return jsonify({"error":"No database"}), 500
        conn = get_db(); cur = conn.cursor()

        cur.execute("SELECT COUNT(*) as n FROM students")
        total_s = cur.fetchone()["n"]

        cur.execute("SELECT COUNT(*) as n FROM events WHERE event_type IN ('login','account_created')")
        total_sess = cur.fetchone()["n"]

        cur.execute("SELECT COUNT(*) as n FROM events WHERE event_type='question_asked'")
        total_q = cur.fetchone()["n"]

        cur.execute("SELECT COUNT(*) as n FROM events WHERE event_type='file_uploaded'")
        total_up = cur.fetchone()["n"]

        # Per-student summary
        cur.execute("""
            SELECT
                s.id, s.first_name, s.last_name, s.email,
                s.classification, s.major,
                to_char(s.created_at, 'Mon DD YYYY') as joined,
                (SELECT COUNT(*) FROM events e WHERE e.student_id=s.id
                 AND e.event_type IN ('login','account_created')) as sessions,
                (SELECT COUNT(*) FROM events e WHERE e.student_id=s.id
                 AND e.event_type='question_asked') as questions,
                (SELECT COUNT(*) FROM events e WHERE e.student_id=s.id
                 AND e.event_type='file_uploaded') as uploads,
                (SELECT COUNT(*) FROM documents d WHERE d.student_id=s.id) as docs
            FROM students s
            ORDER BY s.created_at DESC
        """)
        students = [dict(r) for r in cur.fetchall()]

        # Recent events feed
        cur.execute("""
            SELECT
                e.id, e.event_type, e.payload,
                to_char(e.created_at, 'Mon DD HH24:MI') as ts,
                s.first_name, s.last_name, s.email
            FROM events e
            LEFT JOIN students s ON s.id = e.student_id
            ORDER BY e.created_at DESC
            LIMIT 60
        """)
        recent = []
        for r in cur.fetchall():
            row = dict(r)
            row["payload"] = safe_payload(row.get("payload"))
            recent.append(row)

        cur.execute("SELECT major, COUNT(*) as n FROM students GROUP BY major ORDER BY n DESC")
        by_major = [dict(r) for r in cur.fetchall()]

        cur.execute("SELECT classification, COUNT(*) as n FROM students GROUP BY classification ORDER BY n DESC")
        by_class = [dict(r) for r in cur.fetchall()]

        cur.close(); db_release(conn)
        return jsonify({
            "total_students":  total_s,
            "total_sessions":  total_sess,
            "total_questions": total_q,
            "total_uploads":   total_up,
            "students":        students,
            "recent":          recent,
            "by_major":        by_major,
            "by_class":        by_class
        })
    except Exception as e:
        print(f"analytics_data error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500


@app.route("/analytics-data-full")
def analytics_data_full():
    try:
        s = current_student()
        if not s: return jsonify({"error":"Not logged in"}), 401
        if s["email"].lower() != ADMIN_EMAIL: return jsonify({"error":"Not authorized"}), 403
        if not DB_URL: return jsonify({"error":"No database"}), 500
        conn = get_db(); cur = conn.cursor()

        cur.execute("SELECT COUNT(*) as n FROM students"); total_s = cur.fetchone()["n"]
        cur.execute("SELECT COUNT(*) as n FROM events WHERE event_type IN ('login','account_created')"); total_sess = cur.fetchone()["n"]
        cur.execute("SELECT COUNT(*) as n FROM events WHERE event_type='question_asked'"); total_q = cur.fetchone()["n"]
        cur.execute("SELECT COUNT(*) as n FROM events WHERE event_type='file_uploaded'"); total_up = cur.fetchone()["n"]

        # Per-student summary
        cur.execute("""
            SELECT s.id, s.first_name, s.last_name, s.email, s.classification, s.major,
                   COALESCE(NULLIF(s.university,''),'Not set') as university,
                   to_char(s.created_at,'Mon DD YYYY') as joined,
                   COALESCE(s.is_active, TRUE) as is_active,
                   COALESCE(s.email_verified, FALSE) as email_verified,
                   (SELECT COUNT(*) FROM events e WHERE e.student_id=s.id AND e.event_type IN ('login','account_created')) as sessions,
                   (SELECT COUNT(*) FROM events e WHERE e.student_id=s.id AND e.event_type='question_asked') as questions,
                   (SELECT COUNT(*) FROM events e WHERE e.student_id=s.id AND e.event_type='file_uploaded') as uploads,
                   (SELECT COUNT(*) FROM documents d WHERE d.student_id=s.id) as docs
            FROM students s ORDER BY s.created_at DESC""")
        students = [dict(r) for r in cur.fetchall()]

        # Full questions list (no truncation)
        cur.execute("""
            SELECT e.payload, to_char(e.created_at,'Mon DD HH24:MI') as ts,
                   s.first_name, s.last_name, s.email
            FROM events e LEFT JOIN students s ON s.id=e.student_id
            WHERE e.event_type='question_asked'
            ORDER BY e.created_at DESC LIMIT 200""")
        questions = []
        for r in cur.fetchall():
            row = dict(r)
            p = safe_payload(row.get("payload"))
            questions.append({
                "first_name": row.get("first_name",""),
                "last_name":  row.get("last_name",""),
                "email":      row.get("email",""),
                "question":   p.get("q",""),
                "ts":         row.get("ts","")
            })

        # Paired Q&A conversations
        cur.execute("""
            SELECT e.id, e.event_type, e.payload, e.created_at,
                   to_char(e.created_at,'Mon DD HH24:MI') as ts,
                   s.first_name, s.last_name, s.email, s.id as sid
            FROM events e LEFT JOIN students s ON s.id=e.student_id
            WHERE e.event_type IN ('question_asked','answer_given')
            ORDER BY s.id, e.created_at ASC LIMIT 400""")
        raw_events = [dict(r) for r in cur.fetchall()]
        conversations = []
        i = 0
        while i < len(raw_events):
            ev = raw_events[i]
            p  = safe_payload(ev.get("payload"))
            if ev["event_type"] == "question_asked":
                conv = {
                    "first_name": ev.get("first_name",""),
                    "last_name":  ev.get("last_name",""),
                    "email":      ev.get("email",""),
                    "question":   p.get("q",""),
                    "answer":     "",
                    "ts":         ev.get("ts",""),
                    "sid":        ev.get("sid")
                }
                if i+1 < len(raw_events) and raw_events[i+1]["event_type"] == "answer_given" and raw_events[i+1].get("sid") == ev.get("sid"):
                    ap = safe_payload(raw_events[i+1].get("payload"))
                    conv["answer"] = ap.get("full_answer", "")
                    i += 2
                else:
                    i += 1
                conversations.append(conv)
            else:
                i += 1

        # Recent activity feed (last 100)
        cur.execute("""
            SELECT e.event_type, e.payload, to_char(e.created_at,'Mon DD HH24:MI') as ts,
                   s.first_name, s.last_name, s.email
            FROM events e LEFT JOIN students s ON s.id=e.student_id
            ORDER BY e.created_at DESC LIMIT 100""")
        recent = []
        for r in cur.fetchall():
            row = dict(r)
            row["payload"] = safe_payload(row.get("payload"))
            recent.append(row)

        # By major / classification
        cur.execute("SELECT major, COUNT(*) as n FROM students GROUP BY major ORDER BY n DESC")
        by_major = [dict(r) for r in cur.fetchall()]
        cur.execute("SELECT classification, COUNT(*) as n FROM students GROUP BY classification ORDER BY n DESC")
        by_class = [dict(r) for r in cur.fetchall()]

        # By course (documents)
        cur.execute("SELECT course, COUNT(*) as n FROM documents GROUP BY course ORDER BY n DESC")
        by_course = [dict(r) for r in cur.fetchall()]

        # Daily usage last 7 days
        cur.execute("""
            SELECT to_char(created_at,'Mon DD') as day, COUNT(*) as n
            FROM events
            WHERE created_at >= NOW() - INTERVAL '7 days'
            GROUP BY to_char(created_at,'Mon DD'), DATE(created_at)
            ORDER BY DATE(created_at) ASC""")
        daily = [dict(r) for r in cur.fetchall()]

        # Upcoming deadlines across every student (admin-only visibility into
        # what WINK has extracted from uploaded syllabi/assignment sheets)
        cur.execute("""
            SELECT d.title, d.course, d.due_date, s.first_name, s.last_name
            FROM deadlines d JOIN students s ON s.id = d.student_id
            WHERE d.due_date >= CURRENT_DATE
            ORDER BY d.due_date ASC LIMIT 100""")
        upcoming_deadlines = []
        for r in cur.fetchall():
            row = dict(r)
            row["due_date"] = row["due_date"].isoformat() if row["due_date"] else None
            upcoming_deadlines.append(row)
        cur.execute("SELECT COUNT(*) as n FROM deadlines WHERE due_date >= CURRENT_DATE")
        total_deadlines = cur.fetchone()["n"]

        insights = compute_engagement_insights(cur)

        cur.close(); db_release(conn)
        return jsonify({
            "total_students":  total_s,
            "total_sessions":  total_sess,
            "total_questions": total_q,
            "total_uploads":   total_up,
            "total_deadlines": total_deadlines,
            "students":        students,
            "questions":       questions,
            "conversations":   conversations,
            "recent":          recent,
            "by_major":        by_major,
            "by_class":        by_class,
            "by_course":       by_course,
            "daily":           daily,
            "upcoming_deadlines": upcoming_deadlines,
            **insights
        })
    except Exception as e:
        print(f"analytics_data_full error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500

@app.route("/student-conversations/<int:sid>")
def student_conversations(sid):
    try:
        s = current_student()
        if not s: return jsonify({"error":"Not logged in"}), 401
        if s["email"].lower() != ADMIN_EMAIL: return jsonify({"error":"Not authorized"}), 403
        if not DB_URL: return jsonify({"error":"No database"}), 500
        conn = get_db(); cur = conn.cursor()
        cur.execute("""
            SELECT e.event_type, e.payload, to_char(e.created_at,'Mon DD HH24:MI') as ts
            FROM events e
            WHERE e.student_id=%s AND e.event_type IN ('question_asked','answer_given')
            ORDER BY e.created_at ASC""", (sid,))
        rows = [dict(r) for r in cur.fetchall()]
        cur.close(); db_release(conn)
        conversations = []
        i = 0
        while i < len(rows):
            ev = rows[i]
            p  = safe_payload(ev.get("payload"))
            if ev["event_type"] == "question_asked":
                conv = {"question": p.get("q",""), "answer":"", "ts": ev.get("ts","")}
                if i+1 < len(rows) and rows[i+1]["event_type"] == "answer_given":
                    ap = safe_payload(rows[i+1].get("payload"))
                    conv["answer"] = ap.get("full_answer","")
                    i += 2
                else:
                    i += 1
                conversations.append(conv)
            else:
                i += 1
        return jsonify({"conversations": conversations})
    except Exception as e:
        print(f"student_conversations error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500

@app.route("/toggle-student-active", methods=["POST"])
def toggle_student_active():
    """Admin-only: suspend or reactivate a student account without deleting
    their data. A suspended student can't log in (checked in /login) and any
    existing session is invalidated on their next request (current_student())."""
    try:
        s = current_student()
        if not s: return jsonify({"error":"Not logged in"}), 401
        if s["email"].lower() != ADMIN_EMAIL: return jsonify({"error":"Not authorized"}), 403
        if not DB_URL: return jsonify({"error":"No database"}), 500
        data = request.get_json() or {}
        target_id = data.get("student_id")
        if not target_id:
            return jsonify({"error": "Missing student_id"}), 400
        if str(target_id) == str(s["id"]):
            return jsonify({"error": "You can't suspend your own admin account."}), 400
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT id, is_active FROM students WHERE id=%s", (target_id,))
        target = cur.fetchone()
        if not target:
            cur.close(); db_release(conn)
            return jsonify({"error": "Student not found"}), 404
        new_active = not target["is_active"]
        cur.execute("UPDATE students SET is_active=%s WHERE id=%s", (new_active, target_id))
        conn.commit(); cur.close(); db_release(conn)
        log_event(s["id"], "student_suspended" if not new_active else "student_reactivated", {"target_id": target_id})
        return jsonify({"success": True, "is_active": new_active})
    except Exception as e:
        print(f"toggle_student_active error: {e}"); traceback.print_exc()
        return jsonify({"error": "Something went wrong on our end. Please try again."}), 500


# NOTE: /create-admin, /reset-admin-password, and /admin-check used to live
# here as one-time setup helpers. They've been removed — as written, they
# had no authentication at all, so anyone who found the URL could create or
# take over the admin account with a hardcoded password, or list every
# registered student's name and email. If you need to (re)provision the
# admin account, do it directly in the database instead of via an HTTP route.

@app.route("/health")
def health():
    db_ok = False
    if DB_URL:
        try:
            conn = get_db(); cur = conn.cursor()
            cur.execute("SELECT 1")
            cur.fetchone()
            db_ok = True
            cur.close(); db_release(conn)
        except Exception as e:
            print(f"health check db error: {e}")
    return jsonify({
        "status":  "ok" if (db_ok or not DB_URL) else "degraded",
        "db":      db_ok,
        "api_key": bool(ANTHROPIC_API_KEY),
    })

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT",10000)), debug=False)
