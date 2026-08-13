import json
import logging
import threading
import time

from .. import config
from ..errors import log_error
from ..extensions import get_db
from .retrieval import chunk_text, embed_texts, rank_chunks

logger = logging.getLogger(__name__)

try:
    import pytesseract
    from PIL import Image
    _OCR_AVAILABLE = True
except ImportError:
    _OCR_AVAILABLE = False


def _extract_image_text(filepath, orig_name):
    if not _OCR_AVAILABLE:
        return f"[Image file: {orig_name}]"
    try:
        img = Image.open(filepath)
        text = pytesseract.image_to_string(img).strip()
        if not text:
            return f"[Image file: {orig_name} — no readable text found by OCR]"
        return text
    except Exception as e:
        log_error("services.documents.ocr_failed", e, orig_name=orig_name)
        return f"[Image file: {orig_name}]"


def _zip_bomb_safe(filepath):
    import zipfile
    try:
        with zipfile.ZipFile(filepath) as zf:
            infos = zf.infolist()
            if len(infos) > config.MAX_ZIP_ENTRY_COUNT:
                logger.warning("zip_bomb_check: rejected, %d entries exceeds cap", len(infos))
                return False
            total_uncompressed = sum(i.file_size for i in infos)
            if total_uncompressed > config.MAX_ZIP_UNCOMPRESSED_BYTES:
                logger.warning("zip_bomb_check: rejected, %d uncompressed bytes exceeds cap", total_uncompressed)
                return False
            for i in infos:
                if i.compress_size > 0 and i.file_size / i.compress_size > config.MAX_ZIP_COMPRESSION_RATIO:
                    logger.warning(
                        "zip_bomb_check: rejected, entry %r has a %.0f:1 compression ratio",
                        i.filename, i.file_size / i.compress_size,
                    )
                    return False
            return True
    except zipfile.BadZipFile:
        return False


def extract_text(filepath, orig_name):
    ext = orig_name.rsplit(".", 1)[-1].lower() if "." in orig_name else ""
    text = ""
    if ext in ("docx", "pptx", "xlsx") and not _zip_bomb_safe(filepath):
        logger.warning("extract_text: %s failed the zip-bomb safety check, skipping extraction", orig_name)
        return ""
    try:
        if ext == "txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == "pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(filepath)
                if len(reader.pages) > config.MAX_PDF_PAGES:
                    logger.warning("PDF extract skipped: %d pages exceeds the %d-page cap", len(reader.pages), config.MAX_PDF_PAGES)
                    text = ""
                else:
                    pages = []
                    for i, page in enumerate(reader.pages):
                        try:
                            t = page.extract_text()
                            if t and t.strip():
                                pages.append(f"[Page {i+1}]\n{t.strip()}")
                        except Exception as pe:
                            log_error("services.documents.pdf_page_extract_failed", pe, page=i + 1)
                    text = "\n\n".join(pages)
                    logger.info("PDF extracted %d chars from %d pages", len(text), len(reader.pages))
            except Exception as e:
                log_error("services.documents.pdf_extract_failed", e)
        elif ext == "docx":
            try:
                from docx import Document
                doc = Document(filepath)
                parts = []
                for p in doc.paragraphs:
                    if p.text.strip():
                        parts.append(p.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                text = "\n".join(parts)
                logger.info("DOCX extracted %d chars", len(text))
            except Exception as e:
                log_error("services.documents.docx_extract_failed", e)
        elif ext == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                slides = []
                for i, slide in enumerate(prs.slides):
                    parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text.strip())
                    if parts:
                        slides.append(f"[Slide {i+1}]\n" + "\n".join(parts))
                text = "\n\n".join(slides)
                logger.info("PPTX extracted %d chars", len(text))
            except Exception as e:
                log_error("services.documents.pptx_extract_failed", e)
        elif ext == "xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(filepath, read_only=True, data_only=True)
                sheets = []
                for ws in wb.worksheets:
                    lines = []
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                        if cells:
                            lines.append(" | ".join(cells))
                        if len(lines) >= 2000:  
                            lines.append("[...sheet truncated...]")
                            break
                    if lines:
                        sheets.append(f"[Sheet: {ws.title}]\n" + "\n".join(lines))
                n_sheets = len(wb.worksheets)
                wb.close()
                text = "\n\n".join(sheets)
                logger.info("XLSX extracted %d chars from %d sheet(s)", len(text), n_sheets)
            except Exception as e:
                log_error("services.documents.xlsx_extract_failed", e)
        elif ext in ("jpg", "jpeg", "png"):
            text = _extract_image_text(filepath, orig_name)
        else:
            try:
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            except Exception:
                text = ""
    except Exception as e:
        log_error("services.documents.extract_text_failed", e, orig_name=orig_name); text = ""
    if len(text) > 60000:
        text = text[:60000] + "\n\n[Document truncated at 60,000 characters]"
    return text.strip()


_student_docs_cache = {}
_student_docs_cache_lock = threading.Lock()


def invalidate_student_docs_cache(sid):
    with _student_docs_cache_lock:
        _student_docs_cache.pop(sid, None)


def get_docs(sid):
    if not config.DB_URL:
        return []
    now = time.time()
    with _student_docs_cache_lock:
        cached = _student_docs_cache.get(sid)
        if cached and now - cached[0] < config.STUDENT_DOCS_CACHE_TTL_SECONDS:
            return cached[1]
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT * FROM documents WHERE student_id=%s ORDER BY uploaded_at DESC", (sid,))
        docs = [dict(r) for r in cur.fetchall()]; cur.close()
        with _student_docs_cache_lock:
            _student_docs_cache[sid] = (now, docs)
        return docs
    except Exception as e:
        log_error("services.documents.get_docs", e); return []


def group_docs_by_course(docs):
    groups = {}
    order = []
    for d in docs:
        course = (d.get("course") or "").strip() or "General"
        crn = (d.get("crn") or "").strip()
        label = f"{course} (CRN {crn})" if crn else course
        if label not in groups:
            groups[label] = []
            order.append(label)
        groups[label].append(d)
    order.sort()
    return [(label, groups[label]) for label in order]



def store_document_chunks(document_id, student_id, university, course, orig_name, content):
    if not config.DB_URL or not (content or "").strip():
        return
    header = f"[{orig_name}] ({course})" if course else f"[{orig_name}]"
    chunks = chunk_text(content, header=header)
    if not chunks:
        return
    embeddings = embed_texts(chunks, input_type="document")  
    try:
        conn = get_db(); cur = conn.cursor()
        for i, chunk in enumerate(chunks):
            emb = embeddings[i] if embeddings else None
            cur.execute("""INSERT INTO document_chunks
                           (document_id, student_id, university, chunk_index, content, embedding)
                           VALUES (%s, %s, %s, %s, %s, %s)""",
                        (document_id, student_id, university or "", i, chunk,
                         json.dumps(emb) if emb is not None else None))
        conn.commit(); cur.close()
    except Exception as e:
        log_error("services.documents.store_document_chunks", e)


def get_student_chunks(sid):
    if not config.DB_URL:
        return []
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("""SELECT content, embedding FROM document_chunks
                       WHERE student_id=%s ORDER BY document_id, chunk_index""", (sid,))
        rows = cur.fetchall(); cur.close()
        return [{"content": r["content"], "embedding": json.loads(r["embedding"]) if r["embedding"] else None}
                for r in rows]
    except Exception as e:
        log_error("services.documents.get_student_chunks", e); return []


def get_global_chunks(university):
    if not config.DB_URL:
        return []
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("""SELECT content, embedding FROM document_chunks
                       WHERE student_id IS NULL AND (lower(university)=lower(%s) OR lower(university)='all')
                       ORDER BY document_id, chunk_index""", (university or "",))
        rows = cur.fetchall(); cur.close()
        return [{"content": r["content"], "embedding": json.loads(r["embedding"]) if r["embedding"] else None}
                for r in rows]
    except Exception as e:
        log_error("services.documents.get_global_chunks", e); return []


def build_doc_context(docs, question=None, sid=None):
    if not docs:
        return "\n\nThe student has not uploaded any course documents yet."

    has_content = any((d.get("content") or "").strip() for d in docs)
    if not has_content:
        ctx = f"\n\nThe student has {len(docs)} uploaded file(s) but no text could be extracted. "
        ctx += "Files: " + ", ".join(d["orig_name"] for d in docs)
        return ctx

    total_chars = sum(len((d.get("content") or "")) for d in docs)
    intro = f"\n\n{'='*60}\nSTUDENT'S UPLOADED COURSE DOCUMENTS ({len(docs)} files)\n"

    if total_chars <= config.MAX_DOC_CONTEXT_CHARS:
        ctx = intro
        ctx += ("Every document the student has uploaded is included below in FULL — none have "
                "been shortened or skipped. Never tell the student to re-upload something that "
                "appears here.\n")
        ctx += "Answer questions using the actual content of these documents.\n"
        ctx += "Quote specific text, deadlines, requirements directly from the documents.\n"
        ctx += f"{'='*60}\n\n"
        for i, d in enumerate(docs):
            content = (d.get("content") or "").strip()
            ctx += f"[DOCUMENT {i+1}] {d['orig_name']}\n"
            ctx += f"Course: {d['course']} | Size: {round(d.get('size_bytes',0)/1024,1)} KB\n\n"
            ctx += content if content else "[No text could be extracted from this file]"
            ctx += f"\n\n{'-'*40}\n\n"
        ctx += f"{'='*60}\n"
        return ctx

    if question and sid:
        chunk_rows = get_student_chunks(sid)
        if chunk_rows:
            chunk_texts = [c["content"] for c in chunk_rows]
            chunk_embeddings = [c["embedding"] for c in chunk_rows]
            top = rank_chunks(question, chunk_texts, config.RETRIEVAL_TOP_N_STUDENT_DOCS,
                              chunk_embeddings=chunk_embeddings)
            ctx = intro
            ctx += (f"The student has uploaded more material ({len(docs)} files) than fits in one "
                    f"prompt, so below are the excerpts most relevant to their CURRENT question — "
                    f"not the complete text of every document. Every document they've uploaded is "
                    f"still listed by name so you know it exists; never tell them to re-upload "
                    f"something listed here. If they ask a question that needs a document's FULL "
                    f"text (e.g. 'summarize the whole syllabus'), say you're working from the most "
                    f"relevant excerpts for what they've asked so far and offer to look at a "
                    f"specific section or document if they want more of it.\n"
                    f"Uploaded files: " + ", ".join(d["orig_name"] for d in docs) + "\n")
            ctx += f"{'='*60}\n\n"
            ctx += "\n\n---\n\n".join(top)
            ctx += f"\n\n{'='*60}\n"
            return ctx

    ctx = intro
    ctx += ("Every document the student has uploaded is listed below in full or in part — none "
            "have been skipped. Never tell the student to re-upload something that appears here; "
            "if you only see part of a long one, say you have it but only part of its content, "
            "and offer to look at a specific section if they ask.\n")
    ctx += "Answer questions using the actual content of these documents.\n"
    ctx += "Quote specific text, deadlines, requirements directly from the documents.\n"
    ctx += f"{'='*60}\n\n"
    per_doc_budget = max(config.MAX_DOC_CONTEXT_CHARS // max(len(docs), 1), 2000)
    for i, d in enumerate(docs):
        content = (d.get("content") or "").strip()
        header = f"[DOCUMENT {i+1}] {d['orig_name']}\n"
        header += f"Course: {d['course']} | Size: {round(d.get('size_bytes',0)/1024,1)} KB\n"
        header += f"Content ({len(content)} chars):\n"
        if len(content) > per_doc_budget:
            content = content[:per_doc_budget] + "\n[Shortened here to fit — ask about this document specifically for more of it.]"
        ctx += header
        ctx += content if content else "[No text could be extracted from this file]"
        ctx += f"\n\n{'-'*40}\n\n"
    ctx += f"{'='*60}\n"
    return ctx


_global_docs_cache = {}
_global_docs_cache_lock = threading.Lock()


def invalidate_global_docs_cache(university=None):
    with _global_docs_cache_lock:
        if university is None:
            _global_docs_cache.clear()
        else:
            _global_docs_cache.pop((university or "").lower(), None)
            _global_docs_cache.pop(None, None)  


def get_global_docs(university=None):
    if not config.DB_URL:
        return []
    cache_key = (university or "").lower() or None
    now = time.time()
    with _global_docs_cache_lock:
        cached = _global_docs_cache.get(cache_key)
        if cached and now - cached[0] < config.GLOBAL_DOCS_CACHE_TTL_SECONDS:
            return cached[1]
    try:
        conn = get_db(); cur = conn.cursor()
        if university:
            cur.execute("""SELECT * FROM documents WHERE student_id IS NULL
                           AND (lower(university)=lower(%s) OR lower(university)='all')
                           ORDER BY uploaded_at DESC""",
                        (university,))
        else:
            cur.execute("SELECT * FROM documents WHERE student_id IS NULL ORDER BY uploaded_at DESC")
        docs = [dict(r) for r in cur.fetchall()]; cur.close()
        with _global_docs_cache_lock:
            _global_docs_cache[cache_key] = (now, docs)
        return docs
    except Exception as e:
        log_error("services.documents.get_global_docs", e); return []


def build_global_doc_context(docs, university=None, question=None):
    if not docs:
        return "\n\nNo general reference documents have been added yet."
    label = f" for {university}" if university else ""
    intro = f"\n\n{'='*60}\nGENERAL REFERENCE DOCUMENTS (apply to every student{label}, not just this one)\n"
    footer_note = ("These were uploaded by an administrator and are not visible to the student as "
                   "their own files — don't refer to them as 'your uploaded documents' or mention "
                   "that they were uploaded separately; just use them as background knowledge when "
                   "relevant.\n")
    total_chars = sum(len((d.get("content") or "")) for d in docs)

    if total_chars <= config.MAX_GLOBAL_DOC_CONTEXT_CHARS:
        ctx = intro + footer_note + f"{'='*60}\n\n"
        for i, d in enumerate(docs):
            content = (d.get("content") or "").strip()
            if not content:
                continue
            ctx += f"[REFERENCE {i+1}] {d['orig_name']} ({d.get('course') or 'General'})\n{content}\n\n{'-'*40}\n\n"
        ctx += f"{'='*60}\n"
        return ctx

    if question:
        chunk_rows = get_global_chunks(university)
        if chunk_rows:
            chunk_texts = [c["content"] for c in chunk_rows]
            chunk_embeddings = [c["embedding"] for c in chunk_rows]
            top = rank_chunks(question, chunk_texts, config.RETRIEVAL_TOP_N_GLOBAL_DOCS,
                              chunk_embeddings=chunk_embeddings)
            ctx = intro + footer_note
            ctx += f"{'='*60}\n\n"
            ctx += "\n\n---\n\n".join(top)
            ctx += f"\n\n{'='*60}\n"
            return ctx

    ctx = intro + footer_note + f"{'='*60}\n\n"
    per_doc_budget = max(config.MAX_GLOBAL_DOC_CONTEXT_CHARS // max(len(docs), 1), 2000)
    for i, d in enumerate(docs):
        content = (d.get("content") or "").strip()
        if not content:
            continue
        if len(content) > per_doc_budget:
            content = content[:per_doc_budget] + "\n[Shortened here to fit.]"
        ctx += f"[REFERENCE {i+1}] {d['orig_name']} ({d.get('course') or 'General'})\n{content}\n\n{'-'*40}\n\n"
    ctx += f"{'='*60}\n"
    return ctx
