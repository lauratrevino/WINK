import json
import logging
import time

from psycopg2.extras import execute_values

from .. import config
from ..errors import log_error
from ..extensions import generate_csrf_token, db_cursor
from ..security import login_required, page_login_required, admin_required, file_signature_valid, rate_limited, verified_required
from ..services.analytics import log_event
from ..services.course_colors import ensure_course_colors, purge_course_data_if_gone
from ..services.deadlines import extract_deadlines, insert_deadlines
from ..services.documents import (
    extract_text, get_docs, get_global_docs, group_docs_by_course,
    invalidate_global_docs_cache, invalidate_student_docs_cache,
    store_document_chunks,
)

bp = Blueprint("documents", __name__)
logger = logging.getLogger(__name__)


@bp.route("/course-colors")
@login_required
def course_colors():
    s = g.student
    docs = get_docs(s["id"])
    course_names = sorted({(d.get("course") or "").strip() for d in docs
                            if (d.get("course") or "").strip()}, key=str.lower)
    colors = ensure_course_colors(s["id"], course_names)
    return jsonify({"colors": colors, "courses": course_names})

logger = logging.getLogger(__name__)

try:
    import pytesseract
    from PIL import Image
    _OCR_AVAILABLE = True
except ImportError:
    _OCR_AVAILABLE = False

# Per-file wall-clock budget for text extraction. The existing page/entry/
# size caps (MAX_PDF_PAGES, zip-bomb checks, the 60,000-char output cap)
# bound the RESULT of extraction, but nothing previously bounded the WORK
# done to get there — a file can be well within every one of those caps
# and still be deliberately constructed (e.g. a PDF whose pages are mostly
# dense vector art, or an image sized to make OCR slow) to burn far more
# CPU time than its eventual output size would suggest. This is a
# cooperative, best-effort budget: checked between iterations of the
# page/row/slide loops below, so it can't interrupt a single call that's
# already blocking inside a C extension (a single huge PDF page's
# extract_text(), for instance) — the OCR path gets its own hard,
# subprocess-level timeout instead (see _extract_image_text), since that's
# the one step here that reliably can run away.
_EXTRACTION_TIME_BUDGET_SECONDS = 45
# Tesseract itself gets a harder, separate ceiling — pytesseract runs the
# tesseract binary as a subprocess and this timeout kills that subprocess
# outright, which the cooperative budget above cannot do for a single
# call that's already in flight.
_OCR_TIMEOUT_SECONDS = 20
# A defence against PIL "decompression bomb" images: a small, valid file
# on disk (e.g. a highly-compressed PNG) can decode into a bitmap large
# enough to make OCR (or even just opening the image) consume excessive
# memory and time. This is independent of the 16MB request-size cap,
# which limits the file ON DISK, not the size of the pixel buffer it
# decodes to. 40 million pixels is roughly a 6500x6000 image — generous
# for any real photographed syllabus/whiteboard/handout, well below what
# a crafted file can still decode to within a small file size.
_MAX_OCR_PIXELS = 40_000_000


class _ExtractionBudgetExceeded(Exception):
    """Raised internally to unwind out of a page/row/slide loop once the
    wall-clock budget for this document has been used up. Always caught
    within extract_text() itself — callers just see a (possibly partial)
    text result, the same as any other extraction shortfall."""


def _check_budget(deadline):
    if time.monotonic() > deadline:
        raise _ExtractionBudgetExceeded()


def _extract_image_text(filepath, orig_name):
    if not _OCR_AVAILABLE:
        return f"[Image file: {orig_name}]"
    try:
        img = Image.open(filepath)
        pixels = (img.size[0] or 0) * (img.size[1] or 0)
        if pixels > _MAX_OCR_PIXELS:
            logger.warning(
                "OCR skipped: %s decodes to %d pixels, exceeds the %d-pixel cap",
                orig_name, pixels, _MAX_OCR_PIXELS,
            )
            return f"[Image file: {orig_name} — too large to run OCR on]"
        # timeout= runs tesseract as a subprocess with a hard kill after
        # this many seconds, raising RuntimeError on expiry — this is a
        # real ceiling on OCR time, unlike the cooperative budget checks
        # used elsewhere in this module, which can't interrupt a call
        # that's already blocking inside a C extension.
        text = pytesseract.image_to_string(img, timeout=_OCR_TIMEOUT_SECONDS).strip()
        if not text:
            return f"[Image file: {orig_name} — no readable text found by OCR]"
        return text
    except RuntimeError as e:
        # pytesseract raises plain RuntimeError for its own timeout.
        log_error("services.documents.ocr_timeout", e, orig_name=orig_name)
        return f"[Image file: {orig_name} — OCR took too long and was stopped]"
    except Exception as e:
        log_error("services.documents.ocr_failed", e, orig_name=orig_name)
        return f"[Image file: {orig_name}]"


def _zip_bomb_safe(filepath):
    import zipfile
    try:
        with zipfile.ZipFile(filepath) as zf:
            infos = zf.infolist()
            if len(infos) > config.MAX_ZIP_ENTRY_COUNT:
                logger.warning("zip_bomb_check: rejected, %d entries exceeds cap", len(infos))
                return False
            total_uncompressed = sum(i.file_size for i in infos)
            if total_uncompressed > config.MAX_ZIP_UNCOMPRESSED_BYTES:
                logger.warning("zip_bomb_check: rejected, %d uncompressed bytes exceeds cap", total_uncompressed)
                return False
            for i in infos:
                if i.compress_size > 0 and i.file_size / i.compress_size > config.MAX_ZIP_COMPRESSION_RATIO:
                    logger.warning(
                        "zip_bomb_check: rejected, entry %r has a %.0f:1 compression ratio",
                        i.filename, i.file_size / i.compress_size,
                    )
                    return False
            return True
    except zipfile.BadZipFile:
        return False


def extract_text(filepath, orig_name):
    ext = orig_name.rsplit(".", 1)[-1].lower() if "." in orig_name else ""
    text = ""
    if ext in ("docx", "pptx", "xlsx") and not _zip_bomb_safe(filepath):
        logger.warning("extract_text: %s failed the zip-bomb safety check, skipping extraction", orig_name)
        return ""
    deadline = time.monotonic() + _EXTRACTION_TIME_BUDGET_SECONDS
    budget_hit = False
    try:
        if ext == "txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == "pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(filepath)
                if len(reader.pages) > config.MAX_PDF_PAGES:
                    logger.warning("PDF extract skipped: %d pages exceeds the %d-page cap", len(reader.pages), config.MAX_PDF_PAGES)
                    text = ""
                else:
                    pages = []
                    try:
                        for i, page in enumerate(reader.pages):
                            _check_budget(deadline)
                            try:
                                t = page.extract_text()
                                if t and t.strip():
                                    pages.append(f"[Page {i+1}]\n{t.strip()}")
                            except Exception as pe:
                                log_error("services.documents.pdf_page_extract_failed", pe, page=i + 1)
                    except _ExtractionBudgetExceeded:
                        budget_hit = True
                        pages.append(f"[...extraction stopped early after {_EXTRACTION_TIME_BUDGET_SECONDS}s — remaining pages not processed...]")
                    text = "\n\n".join(pages)
                    logger.info("PDF extracted %d chars from %d pages%s", len(text), len(reader.pages),
                                " (budget exceeded)" if budget_hit else "")
            except Exception as e:
                log_error("services.documents.pdf_extract_failed", e)
        elif ext == "docx":
            try:
                from docx import Document
                doc = Document(filepath)
                parts = []
                try:
                    for i, p in enumerate(doc.paragraphs):
                        if i % 200 == 0:
                            _check_budget(deadline)
                        if p.text.strip():
                            parts.append(p.text.strip())
                    for table in doc.tables:
                        _check_budget(deadline)
                        for row in table.rows:
                            cells = [c.text.strip() for c in row.cells if c.text.strip()]
                            if cells:
                                parts.append(" | ".join(cells))
                except _ExtractionBudgetExceeded:
                    budget_hit = True
                    parts.append(f"[...extraction stopped early after {_EXTRACTION_TIME_BUDGET_SECONDS}s...]")
                text = "\n".join(parts)
                logger.info("DOCX extracted %d chars%s", len(text), " (budget exceeded)" if budget_hit else "")
            except Exception as e:
                log_error("services.documents.docx_extract_failed", e)
        elif ext == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                slides = []
                try:
                    for i, slide in enumerate(prs.slides):
                        _check_budget(deadline)
                        parts = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                parts.append(shape.text.strip())
                        if parts:
                            slides.append(f"[Slide {i+1}]\n" + "\n".join(parts))
                except _ExtractionBudgetExceeded:
                    budget_hit = True
                    slides.append(f"[...extraction stopped early after {_EXTRACTION_TIME_BUDGET_SECONDS}s — remaining slides not processed...]")
                text = "\n\n".join(slides)
                logger.info("PPTX extracted %d chars%s", len(text), " (budget exceeded)" if budget_hit else "")
            except Exception as e:
                log_error("services.documents.pptx_extract_failed", e)
        elif ext == "xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(filepath, read_only=True, data_only=True)
                sheets = []
                try:
                    for ws in wb.worksheets:
                        _check_budget(deadline)
                        lines = []
                        for row in ws.iter_rows(values_only=True):
                            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                            if cells:
                                lines.append(" | ".join(cells))
                            if len(lines) >= 2000:  
                                lines.append("[...sheet truncated...]")
                                break
                        if lines:
                            sheets.append(f"[Sheet: {ws.title}]\n" + "\n".join(lines))
                except _ExtractionBudgetExceeded:
                    budget_hit = True
                    sheets.append(f"[...extraction stopped early after {_EXTRACTION_TIME_BUDGET_SECONDS}s — remaining sheets not processed...]")
                n_sheets = len(wb.worksheets)
                wb.close()
                text = "\n\n".join(sheets)
                logger.info("XLSX extracted %d chars from %d sheet(s)%s", len(text), n_sheets,
                            " (budget exceeded)" if budget_hit else "")
            except Exception as e:
                log_error("services.documents.xlsx_extract_failed", e)
        elif ext in ("jpg", "jpeg", "png"):
            text = _extract_image_text(filepath, orig_name)
        else:
            try:
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            except Exception:
                text = ""
    except Exception as e:
        log_error("services.documents.extract_text_failed", e, orig_name=orig_name); text = ""
    if len(text) > 60000:
        text = text[:60000] + "\n\n[Document truncated at 60,000 characters]"
    return text.strip()


def invalidate_student_docs_cache(sid):
    """No-op, kept so every existing call site (documents.py routes) keeps
    working unchanged. get_docs() no longer caches — see the note there for
    why: this cache used to be a plain in-memory dict, which is only safe
    with a single worker process. This app runs multiple gunicorn workers
    (see Dockerfile's --workers), and each worker has its own separate
    copy of that dict — invalidating one worker's copy after an upload
    never touched any other worker's copy, so a request that happened to
    land on a different worker could keep serving an old or missing
    document list indefinitely. That's a correctness bug, not just a
    performance one, so the cache was removed rather than patched."""
    pass


def get_docs(sid):
    """Always reads from the database — no in-memory caching (see
    invalidate_student_docs_cache() for why). A student's document count
    is small and capped (config.MAX_DOCS_PER_STUDENT), so a single indexed
    SELECT is fast; it isn't worth a caching layer that can silently go
    stale across worker processes."""
    if not config.DB_URL:
        return []
    try:
        with db_cursor() as cur:
            cur.execute("SELECT * FROM documents WHERE student_id=%s ORDER BY uploaded_at DESC", (sid,))
            return [dict(r) for r in cur.fetchall()]
    except Exception as e:
        log_error("services.documents.get_docs", e); return []


def group_docs_by_course(docs):
    groups = {}
    order = []
    for d in docs:
        course = (d.get("course") or "").strip() or "General"
        crn = (d.get("crn") or "").strip()
        label = f"{course} (CRN {crn})" if crn else course
        if label not in groups:
            groups[label] = []
            order.append(label)
        groups[label].append(d)
    order.sort()
    return [(label, groups[label]) for label in order]



def store_document_chunks(document_id, student_id, university, course, orig_name, content):
    if not config.DB_URL or not (content or "").strip():
        return
    header = f"[{orig_name}] ({course})" if course else f"[{orig_name}]"
    chunks = chunk_text(content, header=header)
    if not chunks:
        # content was non-empty but produced zero chunks — treat this the
        # same as a hard failure below, since the document will otherwise
        # look completely normal while being invisible to retrieval.
        _mark_chunking_failed(document_id, "chunk_text() returned no chunks for non-empty content")
        return
    embeddings = embed_texts(chunks, input_type="document")
    try:
        with db_cursor(commit=True) as cur:
            # One batched INSERT instead of one round-trip per chunk — a single
            # large document can chunk into 50+ pieces, and every one of those
            # was a separate network round-trip to the database before this.
            rows = [
                (document_id, student_id, university or "", i, chunk,
                 json.dumps(embeddings[i]) if embeddings and embeddings[i] is not None else None)
                for i, chunk in enumerate(chunks)
            ]
            execute_values(
                cur,
                """INSERT INTO document_chunks
                   (document_id, student_id, university, chunk_index, content, embedding)
                   VALUES %s""",
                rows,
            )
            # Clear any previous failure flag — e.g. a reupload/reprocess that
            # succeeds this time around should self-heal the document's status.
            cur.execute("UPDATE documents SET chunking_failed=FALSE WHERE id=%s", (document_id,))
    except Exception as e:
        log_error("services.documents.store_document_chunks", e, document_id=document_id)
        _mark_chunking_failed(document_id, str(e))


def _mark_chunking_failed(document_id, reason):
    """Best-effort: flip the visible flag so a document with failed chunking
    doesn't silently look identical to a fully-working one. Uses its own
    connection/try because this is already inside (or called right after)
    error handling — a failure here should never mask the original error."""
    if not config.DB_URL or not document_id:
        return
    try:
        with db_cursor(commit=True) as cur:
            cur.execute("UPDATE documents SET chunking_failed=TRUE WHERE id=%s", (document_id,))
    except Exception as e:
        log_error("services.documents._mark_chunking_failed", e, document_id=document_id, reason=reason)


def get_student_chunks(sid, question=None):
    """Returns candidate chunks for a student, narrowed at the database
    level before anything reaches Python.

    Previously this pulled EVERY chunk (and every embedding) belonging to
    the student with no LIMIT — with the 20-document cap and ~60,000-char
    per-document extraction cap, that could mean thousands of chunks and
    their embeddings loaded into application memory for a single retrieval-
    triggered question. When `question` is given, this does a cheap
    server-side keyword pre-filter (ts_rank against the GIN index from
    migration 7c2f19a6d3e1) so only chunks that share vocabulary with the
    question are candidates for the real TF-IDF/neural reranking in
    services/retrieval.py. RETRIEVAL_MAX_CANDIDATE_CHUNKS is a hard
    backstop regardless — including for the no-question fallback path.
    """
    if not config.DB_URL:
        return []
    try:
        with db_cursor() as cur:
            if question and question.strip():
                cur.execute(
                    """SELECT content, embedding FROM document_chunks
                       WHERE student_id=%s
                       ORDER BY ts_rank(to_tsvector('english', content), plainto_tsquery('english', %s)) DESC,
                                document_id, chunk_index
                       LIMIT %s""",
                    (sid, question, config.RETRIEVAL_MAX_CANDIDATE_CHUNKS),
                )
            else:
                cur.execute(
                    """SELECT content, embedding FROM document_chunks
                       WHERE student_id=%s ORDER BY document_id, chunk_index
                       LIMIT %s""",
                    (sid, config.RETRIEVAL_MAX_CANDIDATE_CHUNKS),
                )
            rows = cur.fetchall()
        return [{"content": r["content"], "embedding": json.loads(r["embedding"]) if r["embedding"] else None}
                for r in rows]
    except Exception as e:
        log_error("services.documents.get_student_chunks", e); return []


def get_global_chunks(university, question=None):
    """Same database-level narrowing as get_student_chunks() above, for
    the global/reference document pool — see that function's docstring
    for the full rationale."""
    if not config.DB_URL:
        return []
    try:
        with db_cursor() as cur:
            if question and question.strip():
                cur.execute(
                    """SELECT content, embedding FROM document_chunks
                       WHERE student_id IS NULL AND (lower(university)=lower(%s) OR lower(university)='all')
                       ORDER BY ts_rank(to_tsvector('english', content), plainto_tsquery('english', %s)) DESC,
                                document_id, chunk_index
                       LIMIT %s""",
                    (university or "", question, config.RETRIEVAL_MAX_CANDIDATE_CHUNKS),
                )
            else:
                cur.execute(
                    """SELECT content, embedding FROM document_chunks
                       WHERE student_id IS NULL AND (lower(university)=lower(%s) OR lower(university)='all')
                       ORDER BY document_id, chunk_index
                       LIMIT %s""",
                    (university or "", config.RETRIEVAL_MAX_CANDIDATE_CHUNKS),
                )
            rows = cur.fetchall()
        return [{"content": r["content"], "embedding": json.loads(r["embedding"]) if r["embedding"] else None}
                for r in rows]
    except Exception as e:
        log_error("services.documents.get_global_chunks", e); return []


def build_doc_context(docs, question=None, sid=None):
    if not docs:
        return "\n\nThe student has not uploaded any course documents yet."

    has_content = any((d.get("content") or "").strip() for d in docs)
    if not has_content:
        ctx = f"\n\nThe student has {len(docs)} uploaded file(s) but no text could be extracted. "
        ctx += "Files: " + ", ".join(d["orig_name"] for d in docs)
        return ctx

    total_chars = sum(len((d.get("content") or "")) for d in docs)
    intro = f"\n\n{'='*60}\nSTUDENT'S UPLOADED COURSE DOCUMENTS ({len(docs)} files)\n"

    if total_chars <= config.MAX_DOC_CONTEXT_CHARS:
        ctx = intro
        ctx += ("Every document the student has uploaded is included below in FULL — none have "
                "been shortened or skipped. Never tell the student to re-upload something that "
                "appears here.\n")
        ctx += "Answer questions using the actual content of these documents.\n"
        ctx += "Quote specific text, deadlines, requirements directly from the documents.\n"
        ctx += f"{'='*60}\n\n"
        for i, d in enumerate(docs):
            content = (d.get("content") or "").strip()
            ctx += f"[DOCUMENT {i+1}] {d['orig_name']}\n"
            crn = (d.get("crn") or "").strip()
            course_label = f"{d['course']} (CRN {crn})" if crn else d['course']
            ctx += f"Course: {course_label} | Size: {round(d.get('size_bytes',0)/1024,1)} KB\n\n"
            ctx += content if content else "[No text could be extracted from this file]"
            ctx += f"\n\n{'-'*40}\n\n"
        ctx += f"{'='*60}\n"
        return ctx

    if question and sid:
        chunk_rows = get_student_chunks(sid, question=question)
        if chunk_rows:
            chunk_texts = [c["content"] for c in chunk_rows]
            chunk_embeddings = [c["embedding"] for c in chunk_rows]
            top = rank_chunks(question, chunk_texts, config.RETRIEVAL_TOP_N_STUDENT_DOCS,
                              chunk_embeddings=chunk_embeddings)
            def _labeled(d):
                crn = (d.get("crn") or "").strip()
                return f"{d['orig_name']} ({d['course']}, CRN {crn})" if crn else f"{d['orig_name']} ({d['course']})"
            ctx = intro
            ctx += (f"The student has uploaded more material ({len(docs)} files) than fits in one "
                    f"prompt, so below are the excerpts most relevant to their CURRENT question — "
                    f"not the complete text of every document. Every document they've uploaded is "
                    f"still listed by name so you know it exists; never tell them to re-upload "
                    f"something listed here. If they ask a question that needs a document's FULL "
                    f"text (e.g. 'summarize the whole syllabus'), say you're working from the most "
                    f"relevant excerpts for what they've asked so far and offer to look at a "
                    f"specific section or document if they want more of it.\n"
                    f"Uploaded files: " + ", ".join(_labeled(d) for d in docs) + "\n")
            ctx += f"{'='*60}\n\n"
            ctx += "\n\n---\n\n".join(top)
            ctx += f"\n\n{'='*60}\n"
            return ctx

    ctx = intro
    ctx += ("Every document the student has uploaded is listed below in full or in part — none "
            "have been skipped. Never tell the student to re-upload something that appears here; "
            "if you only see part of a long one, say you have it but only part of its content, "
            "and offer to look at a specific section if they ask.\n")
    ctx += "Answer questions using the actual content of these documents.\n"
    ctx += "Quote specific text, deadlines, requirements directly from the documents.\n"
    ctx += f"{'='*60}\n\n"
    per_doc_budget = max(config.MAX_DOC_CONTEXT_CHARS // max(len(docs), 1), 2000)
    for i, d in enumerate(docs):
        content = (d.get("content") or "").strip()
        header = f"[DOCUMENT {i+1}] {d['orig_name']}\n"
        crn = (d.get("crn") or "").strip()
        course_label = f"{d['course']} (CRN {crn})" if crn else d['course']
        header += f"Course: {course_label} | Size: {round(d.get('size_bytes',0)/1024,1)} KB\n"
        header += f"Content ({len(content)} chars):\n"
        if len(content) > per_doc_budget:
            content = content[:per_doc_budget] + "\n[Shortened here to fit — ask about this document specifically for more of it.]"
        ctx += header
        ctx += content if content else "[No text could be extracted from this file]"
        ctx += f"\n\n{'-'*40}\n\n"
    ctx += f"{'='*60}\n"
    return ctx


def invalidate_global_docs_cache(university=None):
    """No-op, kept so every existing call site keeps working unchanged.
    get_global_docs() no longer caches — same cross-worker staleness
    problem as invalidate_student_docs_cache() above, and this one is
    higher-stakes: it's read on every single chat message, so a stale
    worker could keep feeding students an outdated or removed reference
    document indefinitely."""
    pass


def get_global_docs(university=None):
    """Always reads from the database — no in-memory caching (see
    invalidate_global_docs_cache() for why).

    WARNING: this pulls full `content` (up to 60,000 chars each) for
    EVERY matching global reference document, with no LIMIT — safe to
    call when the caller already knows the result set is small (e.g. an
    admin-facing management page), but NOT safe to call unconditionally
    on a per-chat-message hot path. build_global_doc_context() below used
    to do exactly that on every single chat message regardless of how
    many global documents existed or whether their content would even be
    used — see get_global_docs_total_chars()/get_global_doc_names() for
    the cheap alternatives it uses instead now. This is the same
    architectural problem as the one fixed in get_student_chunks()/
    get_global_chunks() (migration 7c2f19a6d3e1), one layer up: unbounded
    reference material accumulated by admins over a research pilot's
    lifetime, reloaded in full on every message for every student at
    that university, rather than only when actually needed."""
    if not config.DB_URL:
        return []
    try:
        with db_cursor() as cur:
            if university:
                cur.execute("""SELECT * FROM documents WHERE student_id IS NULL
                               AND (lower(university)=lower(%s) OR lower(university)='all')
                               ORDER BY uploaded_at DESC""",
                            (university,))
            else:
                cur.execute("SELECT * FROM documents WHERE student_id IS NULL ORDER BY uploaded_at DESC")
            return [dict(r) for r in cur.fetchall()]
    except Exception as e:
        log_error("services.documents.get_global_docs", e); return []


def get_global_docs_total_chars(university=None):
    """Cheap aggregate — total content length and document count for the
    global reference pool, WITHOUT pulling any content into Python. Lets
    build_global_doc_context() decide which branch it needs (full
    content vs. chunk-based retrieval) before paying for a potentially
    large SELECT * — the decision itself only ever needed the total
    length, never the content itself, until the decision was already made."""
    if not config.DB_URL:
        return 0, 0
    try:
        with db_cursor() as cur:
            if university:
                cur.execute("""SELECT COALESCE(SUM(LENGTH(content)),0) as total_chars, COUNT(*) as n
                               FROM documents WHERE student_id IS NULL
                               AND (lower(university)=lower(%s) OR lower(university)='all')""",
                            (university,))
            else:
                cur.execute("""SELECT COALESCE(SUM(LENGTH(content)),0) as total_chars, COUNT(*) as n
                               FROM documents WHERE student_id IS NULL""")
            row = cur.fetchone()
        return (row["total_chars"], row["n"]) if row else (0, 0)
    except Exception as e:
        log_error("services.documents.get_global_docs_total_chars", e); return 0, 0


def get_global_doc_names(university=None):
    """Just the filenames of the global reference pool — no content — for
    callers (citation verification in chat.py) that need to know what
    documents EXIST without paying for what's IN them. Capped at
    RETRIEVAL_MAX_CANDIDATE_CHUNKS*4 as a defensive backstop; even a
    large research pilot's global reference library shouldn't approach
    that many distinct uploaded files, and this is metadata only (no
    content), so the cap is generous rather than tight."""
    if not config.DB_URL:
        return []
    try:
        limit = config.RETRIEVAL_MAX_CANDIDATE_CHUNKS * 4
        with db_cursor() as cur:
            if university:
                cur.execute("""SELECT orig_name FROM documents WHERE student_id IS NULL
                               AND (lower(university)=lower(%s) OR lower(university)='all')
                               ORDER BY uploaded_at DESC LIMIT %s""",
                            (university, limit))
            else:
                cur.execute("SELECT orig_name FROM documents WHERE student_id IS NULL ORDER BY uploaded_at DESC LIMIT %s",
                            (limit,))
            return [r["orig_name"] for r in cur.fetchall()]
    except Exception as e:
        log_error("services.documents.get_global_doc_names", e); return []


def build_global_doc_context(university=None, question=None):
    """Builds the general-reference-material context for a chat message.

    Deliberately does NOT take a pre-fetched docs list as a parameter
    anymore (it used to) — that meant every caller had to call the
    expensive get_global_docs() first regardless of whether this function
    would even use the content, since the OLD version's total_chars
    check itself required the content to already be in memory to sum its
    length. Checking the cheap aggregate (get_global_docs_total_chars)
    FIRST lets this skip the expensive full fetch entirely on the
    (common, once a research pilot accumulates any real amount of
    reference material) path where retrieval is going to be used anyway.
    """
    total_chars, n_docs = get_global_docs_total_chars(university)
    if n_docs == 0:
        return "\n\nNo general reference documents have been added yet."

    label = f" for {university}" if university else ""
    intro = f"\n\n{'='*60}\nGENERAL REFERENCE DOCUMENTS (apply to every student{label}, not just this one)\n"
    footer_note = ("These were uploaded by an administrator and are not visible to the student as "
                   "their own files — don't refer to them as 'your uploaded documents' or mention "
                   "that they were uploaded separately; just use them as background knowledge when "
                   "relevant.\n")

    if total_chars <= config.MAX_GLOBAL_DOC_CONTEXT_CHARS:
        # Small enough that fetching full content is safe and simple —
        # same full-inclusion behavior as before.
        docs = get_global_docs(university)
        ctx = intro + footer_note + f"{'='*60}\n\n"
        for i, d in enumerate(docs):
            content = (d.get("content") or "").strip()
            if not content:
                continue
            ctx += f"[REFERENCE {i+1}] {d['orig_name']} ({d.get('course') or 'General'})\n{content}\n\n{'-'*40}\n\n"
        ctx += f"{'='*60}\n"
        return ctx

    if question:
        # The retrieval path never needed full document content at all —
        # only the chunk table, which get_global_chunks() already bounds
        # (migration 7c2f19a6d3e1). This is the branch that previously
        # paid for a full get_global_docs() fetch it never used.
        chunk_rows = get_global_chunks(university, question=question)
        if chunk_rows:
            chunk_texts = [c["content"] for c in chunk_rows]
            chunk_embeddings = [c["embedding"] for c in chunk_rows]
            top = rank_chunks(question, chunk_texts, config.RETRIEVAL_TOP_N_GLOBAL_DOCS,
                              chunk_embeddings=chunk_embeddings)
            ctx = intro + footer_note
            ctx += f"{'='*60}\n\n"
            ctx += "\n\n---\n\n".join(top)
            ctx += f"\n\n{'='*60}\n"
            return ctx

    # No question given (or chunk retrieval came up empty) — genuinely
    # needs full content to build the per-document-budget fallback below.
    # This is the one remaining path that pays get_global_docs()'s full
    # cost, but it's the rare/legacy case (chat.py always passes a
    # question), not the common hot path.
    docs = get_global_docs(university)
    ctx = intro + footer_note + f"{'='*60}\n\n"
    per_doc_budget = max(config.MAX_GLOBAL_DOC_CONTEXT_CHARS // max(len(docs), 1), 2000)
    for i, d in enumerate(docs):
        content = (d.get("content") or "").strip()
        if not content:
            continue
        if len(content) > per_doc_budget:
            content = content[:per_doc_budget] + "\n[Shortened here to fit.]"
        ctx += f"[REFERENCE {i+1}] {d['orig_name']} ({d.get('course') or 'General'})\n{content}\n\n{'-'*40}\n\n"
    ctx += f"{'='*60}\n"
    return ctx
